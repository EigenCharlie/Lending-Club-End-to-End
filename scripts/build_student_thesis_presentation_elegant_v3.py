"""Build an expanded elegant student-facing thesis presentation."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(15, 23, 42)
NAVY_2 = RGBColor(30, 41, 59)
BLUE = RGBColor(14, 165, 233)
TEAL = RGBColor(20, 184, 166)
SKY = RGBColor(56, 189, 248)
SLATE = RGBColor(100, 116, 139)
SLATE_LIGHT = RGBColor(226, 232, 240)
TEXT = RGBColor(31, 41, 55)
TEXT_SOFT = RGBColor(71, 85, 105)
ORANGE = RGBColor(249, 115, 22)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(239, 68, 68)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

RESEARCH_QUESTION = (
    "Como mejorar la calibracion y la cuantificacion de la incertidumbre en modelos de "
    "riesgo crediticio (PD, LGD, EAD) mediante tecnicas de prediccion conformal y cual "
    "es su desempeno frente a metodos tradicionales de calibracion."
)

GENERAL_OBJECTIVE = (
    "Implementar y evaluar tecnicas de prediccion conformal para mejorar la calibracion y "
    "cuantificacion de la incertidumbre en modelos de riesgo crediticio (PD, LGD, EAD), "
    "utilizando Lending Club como caso de estudio y comparando su desempeno frente a "
    "metodos tradicionales de calibracion."
)

SPECIFIC_OBJECTIVES = [
    "Realizar una revision sistematica del estado del arte sobre calibracion de modelos, cuantificacion de incertidumbre y prediccion conformal aplicada al riesgo crediticio, identificando variantes como Split Conformal, CQR y Venn-Abers.",
    "Construir y preparar un dataset a partir de los registros publicos de Lending Club (2007-2020), definiendo las variables objetivo PD, LGD y EAD y realizando el preprocesamiento necesario para la experimentacion.",
    "Disenar e implementar experimentos aplicando tecnicas de Conformal Prediction sobre modelos base de PD, LGD y EAD.",
    "Comparar el desempeno de Conformal Prediction con metodos tradicionales de calibracion como Platt Scaling, Isotonic Regression y Bootstrap mediante metricas de calibracion, cobertura, eficiencia y discriminacion.",
    "Evaluar el impacto potencial de la prediccion conformal en provisiones bajo IFRS 9, asignacion de capital regulatorio y metricas de validacion de modelos internos.",
    "Proponer lineamientos practicos y recomendaciones para la adopcion de Conformal Prediction en entornos bancarios reales, considerando implementacion, escalabilidad y gobernanza.",
]

AGENDA_ITEMS = [
    "Introduccion",
    "Justificacion del problema",
    "Estado del arte",
    "Objetivos",
    "Metodologia de solucion propuesta",
    "Resultados",
    "Productos",
    "Conclusiones",
    "Referencias",
]

OBJECTIVE_EVIDENCE_BLOCKS = [
    {
        "objective": SPECIFIC_OBJECTIVES[0],
        "deliverable": (
            "Marco teorico aplicado y tabla del estado del arte con calibracion, Venn-Abers, "
            "Split/Mondrian y CQR."
        ),
        "result": (
            "La tesis no usa Conformal Prediction como palabra de moda: define cuando sirve, "
            "cuando no, y como se conecta con riesgo crediticio."
        ),
    },
    {
        "objective": SPECIFIC_OBJECTIVES[1],
        "deliverable": (
            "Dataset temporal limpio, con splits OOT y variables objetivo listas para PD, LGD y EAD."
        ),
        "result": (
            "Se logro una base reproducible: Lending Club limpio (1.86M prestamos, 110 columnas) "
            "y separado para entrenamiento, calibracion y prueba."
        ),
    },
    {
        "objective": SPECIFIC_OBJECTIVES[2],
        "deliverable": (
            "Bateria comparable de experimentos conformales y criterio explicito para seleccionar la variante final."
        ),
        "result": (
            "El proyecto deja una decision metodologica trazable: que se comparo, que se descarto y "
            "por que la politica final se quedo con Mondrian en PD y adaptativo grade-tiempo en LGD."
        ),
    },
    {
        "objective": SPECIFIC_OBJECTIVES[3],
        "deliverable": (
            "Benchmark OOT de PD y de variantes conformales, con lectura de discriminacion, calibracion, "
            "coverage y eficiencia."
        ),
        "result": (
            "CatBoost calibrado con Isotonic mejora la calidad probabilistica, y el conformal global "
            "demuestra por que una garantia solo promedio no es suficiente."
        ),
    },
    {
        "objective": SPECIFIC_OBJECTIVES[4],
        "deliverable": (
            "Lectura de negocio sobre provisiones: escenarios IFRS9, stages y sensibilidad del ECL."
        ),
        "result": (
            "La tesis llega hasta una conversacion financiera concreta: la incertidumbre mueve provisiones "
            "desde 0.97B en baseline hasta 1.78B en escenario severe."
        ),
    },
    {
        "objective": SPECIFIC_OBJECTIVES[5],
        "deliverable": (
            "Paquete reproducible de adopcion: runbook, checkpoints de gobernanza y narrativa ejecutable."
        ),
        "result": (
            "El entregable final no es solo una corrida tecnica: es una ruta practica para explicar, "
            "monitorear y defender la adopcion de CP en un contexto bancario real."
        ),
    },
]

CRISP_PHASES = [
    {
        "title": "1. Business Understanding",
        "body": "Pasar de ranking a decision bajo incertidumbre con alcance defendible para especializacion.",
    },
    {
        "title": "2. Data Understanding",
        "body": "Validar senal real de riesgo en Lending Club por grade, plazo y faltantes criticos.",
    },
    {
        "title": "3. Data Preparation",
        "body": "Construir splits temporales OOT, contratos de columnas y controles anti-leakage.",
    },
    {
        "title": "4. Modeling",
        "body": "Entrenar PD calibrada y comparar variantes conformales para PD/LGD/EAD.",
    },
    {
        "title": "5. Evaluation",
        "body": "Medir AUC, Brier, coverage, estabilidad temporal, fairness e impacto IFRS9.",
    },
    {
        "title": "6. Deployment",
        "body": "Empaquetar una historia reproducible con runbook y evidencia ejecutable por componente.",
    },
]

REFERENCES = [
    "Basel Committee on Banking Supervision (2006). International Convergence of Capital Measurement and Capital Standards.",
    "IASB (2014). IFRS 9 Financial Instruments.",
    "Vovk, Gammerman y Shafer (2022). Algorithmic Learning in a Random World.",
    "Angelopoulos y Bates (2023). A Gentle Introduction to Conformal Prediction and Distribution-Free Uncertainty Quantification.",
    "Romano, Patterson y Candes (2019). Conformalized Quantile Regression.",
    "Bellotti (2017). Reliable region predictions for automated credit scoring.",
    "Vovk y Petej (2014). Venn-Abers Predictors.",
    "Lending Club (2020). Lending Club Loan Data 2007-2020 Q3.",
]


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--output",
        default="/home/eigenlinux/documentos/Presentacion_tesis_especializacion_elegante.pptx",
    )
    parser.add_argument(
        "--assets-dir",
        default="/home/eigenlinux/documentos/tesis_presentation_assets_elegant_v3",
    )
    return parser.parse_args()


def _load_json(path: Path) -> dict:
    return json.loads(path.read_text())


def _setup_plots() -> None:
    plt.style.use("default")
    plt.rcParams.update(
        {
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "axes.edgecolor": "#CBD5E1",
            "axes.labelcolor": "#334155",
            "xtick.color": "#475569",
            "ytick.color": "#475569",
            "text.color": "#0F172A",
            "font.size": 10,
            "axes.titleweight": "bold",
        }
    )


def _style_axis(ax) -> None:
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def _format_count(value: int) -> str:
    if value >= 1_000_000:
        return f"{value / 1_000_000:.2f}M"
    if value >= 1_000:
        return f"{value / 1_000:.0f}k"
    return str(value)


def _load_split_info(root: Path) -> list[dict]:
    rows = []
    for label, path in [
        ("Train", root / "data/processed/train.parquet"),
        ("Calibration", root / "data/processed/calibration.parquet"),
        ("Test", root / "data/processed/test.parquet"),
    ]:
        df = pd.read_parquet(path, columns=["issue_d", "default_flag"])
        rows.append(
            {
                "label": label,
                "n": int(len(df)),
                "start": pd.to_datetime(df["issue_d"].min()),
                "end": pd.to_datetime(df["issue_d"].max()),
                "default_rate": float(df["default_flag"].mean()),
            }
        )
    return rows


def _save_grade_chart(eda: dict, out: Path) -> None:
    _setup_plots()
    series = pd.Series(eda["default_rate_by_grade"]).astype(float).sort_index()
    colors = ["#D7F0FF", "#B7E6FF", "#8CD7FF", "#62C3F6", "#36ADF0", "#128FDC", "#0B63A9"]
    fig, ax = plt.subplots(figsize=(5.0, 3.0), dpi=220)
    bars = ax.bar(
        series.index, series.values * 100, color=colors, edgecolor="#0F172A", linewidth=0.5
    )
    ax.set_title("Default rate por grade")
    ax.set_ylabel("%")
    ax.set_ylim(0, max(series.values * 100) * 1.22)
    ax.grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(ax)
    for bar, value in zip(bars, series.values, strict=False):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + 0.8,
            f"{value * 100:.1f}%",
            ha="center",
            va="bottom",
            fontsize=7,
        )
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_split_timeline_chart(split_info: list[dict], out: Path) -> None:
    _setup_plots()
    fig, axes = plt.subplots(
        1,
        2,
        figsize=(7.4, 3.0),
        dpi=220,
        gridspec_kw={"width_ratios": [1.7, 1.0]},
    )
    colors = {"Train": "#0EA5E9", "Calibration": "#14B8A6", "Test": "#F97316"}

    ax = axes[0]
    rows = list(reversed(split_info))
    for idx, row in enumerate(rows):
        start = mdates.date2num(row["start"].to_pydatetime())
        end = mdates.date2num(row["end"].to_pydatetime())
        width = end - start
        color = colors[row["label"]]
        ax.barh(
            idx, width, left=start, height=0.42, color=color, edgecolor="#0F172A", linewidth=0.4
        )
        ax.text(
            start + width / 2,
            idx,
            f"{row['label']}  {_format_count(row['n'])}",
            ha="center",
            va="center",
            fontsize=8,
            color="white",
            fontweight="bold",
        )

    ax.set_yticks(range(len(rows)))
    ax.set_yticklabels([row["label"] for row in rows])
    ax.set_title("Split temporal OOT")
    ax.xaxis.set_major_locator(mdates.YearLocator(3))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y"))
    ax.grid(axis="x", linestyle="--", alpha=0.18)
    _style_axis(ax)

    ax2 = axes[1]
    labels = [row["label"] for row in split_info]
    values = [row["default_rate"] * 100 for row in split_info]
    bars = ax2.bar(
        labels,
        values,
        color=[colors[label] for label in labels],
        edgecolor="#0F172A",
        linewidth=0.4,
    )
    ax2.set_title("Default rate por split")
    ax2.set_ylim(0, max(values) * 1.28)
    ax2.grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(ax2)
    for bar, value in zip(bars, values, strict=False):
        ax2.text(
            bar.get_x() + bar.get_width() / 2, value + 0.6, f"{value:.1f}%", ha="center", fontsize=8
        )

    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _draw_diagram_box(
    ax,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    subtitle: str,
    fill: str,
    edge: str,
    title_color: str = "#0F172A",
) -> None:
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle="round,pad=0.012,rounding_size=0.02",
        facecolor=fill,
        edgecolor=edge,
        linewidth=1.2,
        transform=ax.transAxes,
    )
    ax.add_patch(patch)
    ax.text(
        x + w / 2,
        y + h * 0.62,
        title,
        ha="center",
        va="center",
        fontsize=10,
        color=title_color,
        fontweight="bold",
        transform=ax.transAxes,
    )
    ax.text(
        x + w / 2,
        y + h * 0.28,
        subtitle,
        ha="center",
        va="center",
        fontsize=8.2,
        color="#334155",
        transform=ax.transAxes,
    )


def _draw_arrow(
    ax, start: tuple[float, float], end: tuple[float, float], label: str | None = None
) -> None:
    arrow = FancyArrowPatch(
        start,
        end,
        arrowstyle="-|>",
        mutation_scale=12,
        linewidth=1.3,
        color="#94A3B8",
        transform=ax.transAxes,
    )
    ax.add_patch(arrow)
    if label:
        ax.text(
            (start[0] + end[0]) / 2,
            (start[1] + end[1]) / 2 + 0.04,
            label,
            ha="center",
            va="center",
            fontsize=7.8,
            color="#475569",
            transform=ax.transAxes,
        )


def _save_pipeline_diagram(out: Path) -> None:
    _setup_plots()
    fig, ax = plt.subplots(figsize=(10.8, 2.8), dpi=220)
    ax.axis("off")

    _draw_diagram_box(
        ax,
        0.02,
        0.36,
        0.14,
        0.28,
        "Lending Club",
        "Historial publico\ny temporal",
        "#EEF2FF",
        "#94A3B8",
    )
    _draw_diagram_box(
        ax,
        0.20,
        0.36,
        0.16,
        0.28,
        "Limpieza + OOT",
        "Train / Cal /\nTest temporal",
        "#F8FAFC",
        "#94A3B8",
    )
    _draw_diagram_box(
        ax,
        0.40,
        0.36,
        0.16,
        0.28,
        "Feature engineering",
        "WOE/IV y\n60 features",
        "#F8FAFC",
        "#94A3B8",
    )
    _draw_diagram_box(
        ax,
        0.60,
        0.36,
        0.16,
        0.28,
        "PD + calibracion",
        "CatBoost +\nIsotonic",
        "#EFF6FF",
        "#3B82F6",
        title_color="#1D4ED8",
    )
    _draw_diagram_box(
        ax,
        0.80,
        0.36,
        0.16,
        0.28,
        "Conformal Mondrian",
        "Intervalos\npor grade",
        "#EEFDF8",
        "#14B8A6",
        title_color="#0F766E",
    )
    _draw_diagram_box(
        ax,
        0.83,
        0.72,
        0.14,
        0.18,
        "Evaluacion",
        "Coverage +\nbacktest 33 x 7",
        "#ECFDF5",
        "#10B981",
        title_color="#047857",
    )
    _draw_diagram_box(
        ax,
        0.83,
        0.06,
        0.14,
        0.18,
        "Impacto IFRS9",
        "4 escenarios\n3 stages",
        "#FFF7ED",
        "#F97316",
        title_color="#C2410C",
    )

    _draw_arrow(ax, (0.16, 0.50), (0.20, 0.50), "OOT")
    _draw_arrow(ax, (0.36, 0.50), (0.40, 0.50), "WOE / IV")
    _draw_arrow(ax, (0.56, 0.50), (0.60, 0.50), "60 features")
    _draw_arrow(ax, (0.76, 0.50), (0.80, 0.50), "PD calibrada")
    _draw_arrow(ax, (0.96, 0.54), (0.96, 0.72), "Coverage")
    _draw_arrow(ax, (0.96, 0.46), (0.96, 0.24), "Uncertainty")

    ax.text(
        0.01,
        0.95,
        "Pipeline canonico del proyecto",
        fontsize=12,
        fontweight="bold",
        color="#0F172A",
        transform=ax.transAxes,
    )
    ax.text(
        0.01,
        0.89,
        "Del dato crudo a la lectura de negocio, con trazabilidad metodologica en cada salto.",
        fontsize=8.5,
        color="#475569",
        transform=ax.transAxes,
    )

    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_conformal_schematic(out: Path) -> None:
    _setup_plots()
    fig, axes = plt.subplots(1, 3, figsize=(9.1, 3.3), dpi=220)

    preds = [0.03, 0.05, 0.08, 0.11, 0.12, 0.18, 0.24, 0.29]
    y_vals = list(range(1, len(preds) + 1))
    axes[0].scatter(preds, y_vals, s=48, color="#0EA5E9", edgecolor="#0F172A", linewidth=0.4)
    axes[0].scatter(
        [0.12], [5], s=90, color="#F97316", edgecolor="#0F172A", linewidth=0.6, zorder=3
    )
    axes[0].annotate(
        "nuevo cliente\nPD = 12%",
        xy=(0.12, 5),
        xytext=(0.18, 6.3),
        fontsize=8,
        arrowprops={"arrowstyle": "->", "color": "#475569"},
    )
    axes[0].set_title("1. Modelo puntual")
    axes[0].set_xlabel("PD estimada")
    axes[0].set_ylabel("Observacion")
    axes[0].set_xlim(0, 0.34)
    axes[0].set_ylim(0.5, 8.7)
    axes[0].grid(alpha=0.15)
    _style_axis(axes[0])

    scores = [0.01, 0.02, 0.02, 0.03, 0.04, 0.05, 0.05, 0.06, 0.07, 0.08, 0.10, 0.12, 0.14]
    q90 = pd.Series(scores).quantile(0.90)
    axes[1].hist(scores, bins=7, color="#DBEAFE", edgecolor="#1D4ED8", linewidth=0.8)
    axes[1].axvline(q90, color="#F97316", linestyle="--", linewidth=1.6)
    axes[1].text(
        q90 + 0.004, 2.15, f"q90 = {q90:.2f}", color="#C2410C", fontsize=8, fontweight="bold"
    )
    axes[1].set_title("2. Set de calibracion")
    axes[1].set_xlabel("Score de no conformidad")
    axes[1].set_ylabel("Frecuencia")
    axes[1].grid(axis="y", alpha=0.15)
    _style_axis(axes[1])

    labels = ["Grade A", "Grade C", "Grade G"]
    lows = [0.04, 0.08, 0.05]
    mids = [0.06, 0.12, 0.15]
    highs = [0.08, 0.17, 0.28]
    colors = ["#0EA5E9", "#14B8A6", "#F97316"]
    axes[2].set_title("3. Intervalo conformal")
    for idx, (label, low, mid, high, color) in enumerate(
        zip(labels, lows, mids, highs, colors, strict=False), start=1
    ):
        axes[2].hlines(idx, low, high, color=color, linewidth=5.0)
        axes[2].scatter([mid], [idx], s=46, color="#0F172A", zorder=3)
        axes[2].text(high + 0.01, idx, f"{label}: [{low:.0%}, {high:.0%}]", va="center", fontsize=8)
    axes[2].set_xlabel("Rango plausible de PD")
    axes[2].set_xlim(0, 0.36)
    axes[2].set_ylim(0.5, 3.7)
    axes[2].set_yticks([])
    axes[2].grid(axis="x", alpha=0.15)
    _style_axis(axes[2])

    fig.suptitle(
        "Conformal Prediction agrega un rango confiable alrededor de la prediccion puntual",
        fontsize=12,
        fontweight="bold",
        y=1.03,
    )
    fig.text(
        0.5,
        -0.02,
        "Esquema conceptual: se usa un set de calibracion para transformar errores pasados en una garantia de cobertura futura.",
        ha="center",
        fontsize=8.5,
        color="#475569",
    )
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_pd_chart(roc: pd.DataFrame, cal: pd.DataFrame, out: Path) -> None:
    _setup_plots()
    fig, axes = plt.subplots(1, 2, figsize=(8.6, 3.6), dpi=220)
    roc_specs = {
        "logreg": ("Logistic baseline", "#94A3B8"),
        "catboost_calibrated": ("CatBoost calibrado", "#0EA5E9"),
    }
    for model, (label, color) in roc_specs.items():
        subset = roc.loc[roc["model"] == model]
        axes[0].plot(subset["fpr"], subset["tpr"], color=color, linewidth=2.5, label=label)
    axes[0].plot([0, 1], [0, 1], linestyle="--", color="#CBD5E1", linewidth=1.2)
    axes[0].set_title("ROC")
    axes[0].set_xlabel("FPR")
    axes[0].set_ylabel("TPR")
    axes[0].grid(alpha=0.18)
    axes[0].legend(frameon=False, fontsize=8)

    cal_specs = {
        "catboost_uncalibrated": ("Sin calibrar", "#F59E0B"),
        "catboost_calibrated": ("Isotonic", "#0EA5E9"),
    }
    for model, (label, color) in cal_specs.items():
        subset = cal.loc[cal["model"] == model]
        axes[1].plot(
            subset["predicted_prob"],
            subset["observed_freq"],
            color=color,
            linewidth=2.4,
            marker="o",
            markersize=3.2,
            label=label,
        )
    axes[1].plot([0, 1], [0, 1], linestyle="--", color="#CBD5E1", linewidth=1.2)
    axes[1].set_title("Curva de calibracion")
    axes[1].set_xlabel("Probabilidad predicha")
    axes[1].set_ylabel("Frecuencia observada")
    axes[1].set_xlim(0, 1)
    axes[1].set_ylim(0, 1)
    axes[1].grid(alpha=0.18)
    axes[1].legend(frameon=False, fontsize=8)

    for ax in axes:
        _style_axis(ax)

    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_coverage_target_chart(monthly: pd.DataFrame, out: Path) -> None:
    _setup_plots()
    fig, ax = plt.subplots(figsize=(7.7, 3.3), dpi=220)
    months = pd.to_datetime(monthly["month"])
    ax.plot(
        months, monthly["coverage_90"] * 100, color="#0EA5E9", linewidth=2.2, label="Coverage 90%"
    )
    ax.plot(
        months, monthly["coverage_95"] * 100, color="#14B8A6", linewidth=2.2, label="Coverage 95%"
    )
    ax.axhline(90, linestyle="--", color="#F97316", linewidth=1.3, label="Meta 90%")
    ax.axhline(95, linestyle="--", color="#FB7185", linewidth=1.3, label="Meta 95%")
    ax.set_title("Backtest temporal de cobertura")
    ax.set_ylabel("Cobertura (%)")
    ax.set_ylim(88, 97.2)
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=4))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m"))
    ax.tick_params(axis="x", rotation=35, labelsize=7)
    ax.grid(axis="y", linestyle="--", alpha=0.18)
    ax.legend(frameon=False, ncol=2, fontsize=8, loc="lower left")
    _style_axis(ax)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_group_coverage_chart(group_metrics: pd.DataFrame, out: Path) -> None:
    _setup_plots()
    df = group_metrics.rename(columns={"group": "grade"}).copy()
    fig, axes = plt.subplots(1, 2, figsize=(8.3, 3.5), dpi=220)
    colors = ["#BAE6FD", "#7DD3FC", "#38BDF8", "#0EA5E9", "#14B8A6", "#0F766E", "#F97316"]

    bars = axes[0].bar(
        df["grade"], df["coverage_90"] * 100, color=colors, edgecolor="#0F172A", linewidth=0.4
    )
    axes[0].axhline(90, linestyle="--", color="#F97316", linewidth=1.3)
    axes[0].set_title("Cobertura empirica por grade")
    axes[0].set_ylabel("%")
    axes[0].set_ylim(84, 95)
    axes[0].grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(axes[0])
    for bar, cov, n in zip(bars, df["coverage_90"] * 100, df["n"], strict=False):
        axes[0].text(
            bar.get_x() + bar.get_width() / 2, cov + 0.3, f"{cov:.1f}", ha="center", fontsize=7
        )
        axes[0].text(
            bar.get_x() + bar.get_width() / 2,
            84.15,
            f"n={_format_count(int(n))}",
            ha="center",
            va="bottom",
            fontsize=6.8,
            color="#64748B",
        )

    bars2 = axes[1].bar(
        df["grade"], df["avg_width_90"], color=colors, edgecolor="#0F172A", linewidth=0.4
    )
    axes[1].set_title("Ancho promedio del intervalo (90%)")
    axes[1].set_ylabel("Width")
    axes[1].grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(axes[1])
    for bar, value in zip(bars2, df["avg_width_90"], strict=False):
        axes[1].text(
            bar.get_x() + bar.get_width() / 2, value + 0.02, f"{value:.2f}", ha="center", fontsize=7
        )

    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_global_vs_mondrian_chart(
    benchmark_by_group: pd.DataFrame, benchmark: pd.DataFrame, out: Path
) -> None:
    _setup_plots()
    fig, axes = plt.subplots(
        1, 2, figsize=(8.7, 3.6), dpi=220, gridspec_kw={"width_ratios": [1.4, 1.0]}
    )

    keep = benchmark_by_group.loc[
        benchmark_by_group["variant"].isin(["global_split", "mondrian_selected_cfg"])
    ].copy()
    pivot = keep.pivot(index="group", columns="variant", values="coverage").reindex(list("ABCDEFG"))
    x_vals = list(range(len(pivot.index)))
    axes[0].bar(
        [x - 0.18 for x in x_vals],
        pivot["global_split"] * 100,
        width=0.34,
        color="#CBD5E1",
        edgecolor="#0F172A",
        linewidth=0.4,
        label="Global split",
    )
    axes[0].bar(
        [x + 0.18 for x in x_vals],
        pivot["mondrian_selected_cfg"] * 100,
        width=0.34,
        color="#0EA5E9",
        edgecolor="#0F172A",
        linewidth=0.4,
        label="Mondrian",
    )
    axes[0].axhline(90, linestyle="--", color="#F97316", linewidth=1.3)
    axes[0].set_xticks(x_vals)
    axes[0].set_xticklabels(list(pivot.index))
    axes[0].set_ylim(50, 97)
    axes[0].set_ylabel("%")
    axes[0].set_title("Cobertura por grade: global vs Mondrian")
    axes[0].grid(axis="y", linestyle="--", alpha=0.18)
    axes[0].legend(frameon=False, fontsize=8, loc="upper left")
    _style_axis(axes[0])

    name_map = {
        "global_split": "Global",
        "mondrian_scaled": "Mondrian scaled",
        "mondrian_unscaled": "Mondrian raw",
        "mondrian_selected_cfg": "Mondrian sel.",
    }
    for _, row in benchmark.iterrows():
        x = float(row["avg_width"])
        y = float(row["min_group_coverage"]) * 100
        label = name_map.get(str(row["variant"]), str(row["variant"]))
        selected = str(row["variant"]) == "mondrian_selected_cfg"
        axes[1].scatter(
            [x],
            [y],
            s=150 if selected else 90,
            color="#0EA5E9" if selected else "#94A3B8",
            edgecolors="#0F172A",
            linewidths=0.6,
            zorder=3,
        )
        axes[1].text(x + 0.01, y + 0.6, label, fontsize=7.5)
    axes[1].axhline(90, linestyle="--", color="#F97316", linewidth=1.3)
    axes[1].set_xlabel("Ancho promedio")
    axes[1].set_ylabel("Min group coverage (%)")
    axes[1].set_title("Trade-off cobertura - eficiencia")
    axes[1].grid(alpha=0.18)
    _style_axis(axes[1])

    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_lgd_ead_chart(lgd_benchmark: pd.DataFrame, status: dict, out: Path) -> None:
    _setup_plots()
    fig, axes = plt.subplots(1, 2, figsize=(8.5, 3.5), dpi=220)

    short_names = {
        "two_stage_split": "Two-stage",
        "direct_split": "Direct split",
        "direct_cqr": "Direct CQR",
        "direct_adaptive_grade_temporal": "Adaptive\ngrade-time",
    }
    x_labels = [short_names.get(name, name) for name in lgd_benchmark["variant"]]
    bars = axes[0].bar(
        x_labels,
        lgd_benchmark["coverage_90"] * 100,
        color=["#CBD5E1", "#CBD5E1", "#FDE68A", "#14B8A6"],
        edgecolor="#0F172A",
        linewidth=0.4,
    )
    axes[0].axhline(90, linestyle="--", color="#F97316", linewidth=1.3)
    axes[0].set_title("LGD: benchmark de variantes (90%)")
    axes[0].set_ylabel("%")
    axes[0].set_ylim(70, 96)
    axes[0].grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(axes[0])
    for bar, value in zip(bars, lgd_benchmark["coverage_90"] * 100, strict=False):
        axes[0].text(
            bar.get_x() + bar.get_width() / 2, value + 0.5, f"{value:.1f}", ha="center", fontsize=7
        )

    lgd = status["lgd"]["conformal"]
    ead = status["ead"]["conformal"]
    labels = ["LGD 90", "LGD 95", "EAD 90", "EAD 95"]
    values = [
        lgd["metrics_90"]["empirical_coverage"] * 100,
        lgd["metrics_95"]["empirical_coverage"] * 100,
        ead["metrics_90"]["empirical_coverage"] * 100,
        ead["metrics_95"]["empirical_coverage"] * 100,
    ]
    targets = [90, 95, 90, 95]
    bars2 = axes[1].bar(
        labels,
        values,
        color=["#0EA5E9", "#38BDF8", "#14B8A6", "#5EEAD4"],
        edgecolor="#0F172A",
        linewidth=0.4,
    )
    axes[1].scatter(labels, targets, color="#F97316", marker="_", s=400, linewidths=2.4, zorder=5)
    axes[1].set_title("Cobertura final seleccionada")
    axes[1].set_ylabel("%")
    axes[1].set_ylim(86, 99)
    axes[1].grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(axes[1])
    for bar, value in zip(bars2, values, strict=False):
        axes[1].text(
            bar.get_x() + bar.get_width() / 2, value + 0.25, f"{value:.2f}", ha="center", fontsize=7
        )

    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _save_ifrs9_chart(ifrs9: pd.DataFrame, out: Path) -> None:
    _setup_plots()
    order = ["baseline", "mild_stress", "adverse", "severe"]
    df = ifrs9.set_index("scenario").loc[order].reset_index()
    values = df["total_ecl"].astype(float) / 1e9
    colors = ["#BAE6FD", "#7DD3FC", "#38BDF8", "#F97316"]
    fig, ax = plt.subplots(figsize=(6.0, 3.5), dpi=220)
    bars = ax.bar(df["scenario"], values, color=colors, edgecolor="#0F172A", linewidth=0.5)
    ax.set_title("ECL total por escenario")
    ax.set_ylabel("Miles de millones")
    ax.grid(axis="y", linestyle="--", alpha=0.18)
    _style_axis(ax)
    for bar, value in zip(bars, values, strict=False):
        ax.text(
            bar.get_x() + bar.get_width() / 2, value + 0.03, f"{value:.2f}", ha="center", fontsize=8
        )
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def _generate_assets(root: Path, assets_dir: Path) -> tuple[dict[str, Path], dict]:
    assets_dir.mkdir(parents=True, exist_ok=True)
    charts_dir = assets_dir / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)

    eda = _load_json(root / "data/processed/eda_summary.json")
    shapes = _load_json(root / "data/processed/dataset_shapes_summary.json")
    split_info = _load_split_info(root)
    roc = pd.read_parquet(root / "data/processed/roc_curve_data.parquet")
    cal = pd.read_parquet(root / "data/processed/calibration_curve_data.parquet")
    monthly = pd.read_parquet(root / "data/processed/conformal_backtest_monthly.parquet")
    group_metrics = pd.read_parquet(
        root / "data/processed/conformal_group_metrics_mondrian.parquet"
    )
    benchmark = pd.read_parquet(root / "data/processed/conformal_variant_benchmark.parquet")
    benchmark_by_group = pd.read_parquet(
        root / "data/processed/conformal_variant_benchmark_by_group.parquet"
    )
    lgd_benchmark = pd.read_parquet(root / "data/processed/lgd_variant_benchmark.parquet")
    lgd_ead = _load_json(root / "models/conformal_lgd_ead_status.json")
    ifrs9 = pd.read_parquet(root / "data/processed/ifrs9_scenario_summary.parquet")

    paths = {
        "grade": charts_dir / "grade_default.png",
        "splits": charts_dir / "split_timeline.png",
        "pipeline": charts_dir / "pipeline.png",
        "conformal": charts_dir / "conformal_schematic.png",
        "pd": charts_dir / "pd_diagnostics.png",
        "coverage": charts_dir / "coverage_backtest.png",
        "group_coverage": charts_dir / "group_coverage.png",
        "mondrian": charts_dir / "mondrian_vs_global.png",
        "lgd_ead": charts_dir / "lgd_ead.png",
        "ifrs9": charts_dir / "ifrs9.png",
    }

    _save_grade_chart(eda, paths["grade"])
    _save_split_timeline_chart(split_info, paths["splits"])
    _save_pipeline_diagram(paths["pipeline"])
    _save_conformal_schematic(paths["conformal"])
    _save_pd_chart(roc, cal, paths["pd"])
    _save_coverage_target_chart(monthly, paths["coverage"])
    _save_group_coverage_chart(group_metrics, paths["group_coverage"])
    _save_global_vs_mondrian_chart(benchmark_by_group, benchmark, paths["mondrian"])
    _save_lgd_ead_chart(lgd_benchmark, lgd_ead, paths["lgd_ead"])
    _save_ifrs9_chart(ifrs9, paths["ifrs9"])

    data = {
        "eda": eda,
        "shapes": shapes,
        "split_info": split_info,
        "monthly": monthly,
        "group_metrics": group_metrics,
        "benchmark": benchmark,
        "lgd_benchmark": lgd_benchmark,
        "lgd_ead": lgd_ead,
        "ifrs9": ifrs9,
    }
    return paths, data


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_face: str = "Aptos",
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_face
    p.alignment = align


def _add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: RGBColor = TEXT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"- {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Aptos"
        p.space_after = Pt(7)


def _add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = SLATE_LIGHT,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)


def _add_metric(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    value: str,
    *,
    accent: RGBColor = BLUE,
) -> None:
    _add_card(slide, x, y, w, h, fill=WHITE, line=SLATE_LIGHT)
    accent_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()
    _add_text(slide, label, x + 0.16, y + 0.1, w - 0.25, 0.22, size=11, color=TEXT_SOFT, bold=True)
    _add_text(slide, value, x + 0.16, y + 0.36, w - 0.25, h - 0.4, size=20, color=NAVY, bold=True)


def _decorate_content_slide(slide, section: str, slide_no: int, title: str) -> None:
    _set_bg(slide, WHITE)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(0.3), Inches(0.18), Inches(0.58)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    _add_text(
        slide,
        section.upper(),
        11.12,
        0.32,
        1.35,
        0.18,
        size=9,
        color=TEXT_SOFT,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    _add_text(
        slide,
        title,
        0.9,
        0.24,
        10.95,
        0.88,
        size=22,
        color=NAVY,
        bold=True,
        font_face="Aptos Display",
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(1.18), Inches(11.2), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SLATE_LIGHT
    line.line.fill.background()
    dot = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(12.28), Inches(0.34), Inches(0.18), Inches(0.18)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()
    _add_text(slide, str(slide_no), 12.52, 0.28, 0.25, 0.2, size=10, color=TEXT_SOFT, bold=True)


def _add_footer(slide, text: str) -> None:
    _add_text(slide, text, 0.9, 7.04, 11.7, 0.18, size=10, color=SLATE)


def _add_picture(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(
        path.as_posix(), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )


def _add_objective_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    objective: str,
    deliverable: str,
    result: str,
    accent: RGBColor,
) -> None:
    _add_card(slide, x, y, w, h, fill=WHITE, line=SLATE_LIGHT)
    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    _add_text(slide, "Objetivo", x + 0.16, y + 0.10, 1.2, 0.18, size=10, color=TEXT_SOFT, bold=True)
    _add_text(slide, objective, x + 0.16, y + 0.28, w - 0.32, 0.55, size=13, color=NAVY, bold=True)
    _add_text(
        slide, "Entregable", x + 0.16, y + 0.90, 1.2, 0.18, size=10, color=TEXT_SOFT, bold=True
    )
    _add_text(slide, deliverable, x + 0.16, y + 1.07, w - 0.32, 0.36, size=11, color=TEXT)
    _add_text(
        slide, "Que demostro", x + 0.16, y + 1.48, 1.35, 0.18, size=10, color=TEXT_SOFT, bold=True
    )
    _add_text(slide, result, x + 0.16, y + 1.65, w - 0.32, 0.45, size=11, color=TEXT_SOFT)


def _add_phase_card(
    slide, x: float, y: float, w: float, h: float, title: str, body: str, accent: RGBColor
) -> None:
    _add_card(slide, x, y, w, h, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, title, x + 0.14, y + 0.14, w - 0.24, 0.32, size=12, color=accent, bold=True)
    _add_text(slide, body, x + 0.14, y + 0.52, w - 0.24, h - 0.64, size=11, color=TEXT)


def _title_slide(slide) -> None:
    _set_bg(slide, NAVY)
    for x, y, size, color, transparency in [
        (10.9, 0.28, 1.5, SKY, 0.18),
        (11.8, 1.08, 1.0, BLUE, 0.15),
        (10.1, 1.7, 0.72, TEAL, 0.15),
        (11.25, 2.15, 0.58, WHITE, 0.78),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.fill.transparency = transparency
        shape.line.fill.background()

    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(10.1), Inches(4.55), Inches(3.0), Inches(2.2)
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = BLUE
    tri.fill.transparency = 0.1
    tri.line.fill.background()

    _add_text(
        slide,
        "Conformal Prediction para la calibracion y\ncuantificacion de incertidumbre en riesgo crediticio",
        0.88,
        1.34,
        8.2,
        1.65,
        size=28,
        color=WHITE,
        bold=True,
        font_face="Aptos Display",
    )
    _add_text(
        slide,
        "Presentacion de tesis de especializacion para audiencia no experta",
        0.9,
        3.18,
        6.4,
        0.3,
        size=16,
        color=SKY,
        bold=True,
    )
    _add_text(
        slide,
        "Carlos Alfredo Vergara Rojas\nEspecializacion en Analitica y Ciencia de Datos Aplicada - UTP",
        0.9,
        3.62,
        6.0,
        0.62,
        size=14,
        color=WHITE,
    )
    _add_text(
        slide,
        "Idea fuerza: una probabilidad puntual ordena clientes; una prediccion con incertidumbre permite decidir con prudencia.",
        0.9,
        4.72,
        7.45,
        0.45,
        size=14,
        color=SLATE_LIGHT,
    )
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(5.52), Inches(1.8), Inches(0.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def _agenda_slide(slide, slide_no: int) -> None:
    _set_bg(slide, NAVY)
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.78),
        Inches(0.76),
        Inches(11.75),
        Inches(5.95),
    )
    frame.fill.solid()
    frame.fill.transparency = 1.0
    frame.line.color.rgb = RGBColor(96, 165, 250)
    frame.line.width = Pt(2.2)

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.22),
        Inches(0.46),
        Inches(5.55),
        Inches(0.96),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = SKY
    header.line.fill.background()
    _add_text(
        slide,
        "CONTENIDO",
        2.07,
        0.62,
        3.9,
        0.4,
        size=30,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_face="Aptos Display",
    )

    y = 1.86
    for idx, item in enumerate(AGENDA_ITEMS, start=1):
        _add_text(
            slide,
            str(idx),
            1.55,
            y,
            0.42,
            0.28,
            size=26,
            color=SKY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text(slide, item, 2.18, y - 0.03, 7.2, 0.34, size=24, color=WHITE, bold=False)
        y += 0.64

    _add_text(slide, str(slide_no), 12.45, 0.28, 0.3, 0.2, size=10, color=SLATE_LIGHT, bold=True)
    _add_footer(
        slide,
        "La agenda replica el orden recomendado en el feedback: mas espacio para la propuesta antes de entrar a resultados.",
    )


def _slide_justification(slide, slide_no: int) -> None:
    _decorate_content_slide(
        slide,
        "justificacion",
        slide_no,
        "El problema es relevante porque las predicciones puntuales no bastan para un contexto regulado",
    )
    cards = [
        (
            0.95,
            1.52,
            "Dependencia regulatoria",
            "PD, LGD y EAD alimentan provisiones IFRS 9, capital regulatorio y validacion interna.",
            BLUE,
        ),
        (
            4.10,
            1.52,
            "Limitacion tecnica actual",
            "Muchos modelos entregan una prediccion puntual, pero no una medida fiable de incertidumbre.",
            ORANGE,
        ),
        (
            7.25,
            1.52,
            "Consecuencia operativa",
            "Eso puede producir perdidas inesperadas, exceso de provisiones y menor confianza de auditores.",
            RED,
        ),
    ]
    for x, y, title, body, accent in cards:
        _add_phase_card(slide, x, y, 2.78, 2.10, title, body, accent)

    _add_card(
        slide, 10.40, 1.52, 1.80, 2.10, fill=RGBColor(240, 253, 250), line=RGBColor(94, 234, 212)
    )
    _add_text(slide, "Oportunidad", 10.62, 1.82, 1.2, 0.2, size=13, color=TEAL, bold=True)
    _add_text(
        slide,
        "Conformal Prediction agrega cobertura estadistica y mejora la auditabilidad del modelo.",
        10.62,
        2.20,
        1.3,
        0.92,
        size=11,
        color=TEXT,
    )

    _add_card(
        slide, 0.95, 4.10, 11.25, 2.08, fill=RGBColor(239, 246, 255), line=RGBColor(147, 197, 253)
    )
    _add_text(
        slide, "Pregunta de investigacion", 1.20, 4.38, 2.2, 0.2, size=14, color=BLUE, bold=True
    )
    _add_text(
        slide,
        RESEARCH_QUESTION,
        1.20,
        4.70,
        10.6,
        0.92,
        size=18,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_footer(
        slide,
        "Justificacion basada en la propuesta: necesidad academica, regulatoria y practica de cuantificar incertidumbre.",
    )


def _slide_state_of_art(slide, slide_no: int, chart: Path) -> None:
    _decorate_content_slide(
        slide,
        "estado del arte",
        slide_no,
        "El estado del arte separa calibracion probabilistica de cuantificacion formal de incertidumbre",
    )
    _add_card(slide, 0.95, 1.50, 4.25, 4.72, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 1.14, 1.84, 3.86, 3.98)

    method_cards = [
        (
            5.50,
            1.50,
            "Platt / Isotonic",
            "Mejoran calibracion de la PD, pero no entregan cobertura intervalar.",
            BLUE,
        ),
        (
            8.95,
            1.50,
            "Venn-Abers",
            "Produce probabilidades mas robustas, pero no sustituye una politica completa de coverage.",
            TEAL,
        ),
        (
            5.50,
            3.72,
            "Split / Mondrian",
            "Conformal agrega cobertura distribution-free; Mondrian la refina por segmento.",
            ORANGE,
        ),
        (
            8.95,
            3.72,
            "CQR",
            "Es relevante para LGD y EAD porque adapta intervalos a heteroscedasticidad.",
            GREEN,
        ),
    ]
    for x, y, title, body, accent in method_cards:
        _add_phase_card(slide, x, y, 2.95, 1.86, title, body, accent)

    _add_footer(
        slide,
        "Lectura del estado del arte: la brecha no es solo metodologica; tambien es de integracion PD/LGD/EAD con trazabilidad regulatoria.",
    )


def _slide_objectives_overview(slide, slide_no: int) -> None:
    _decorate_content_slide(
        slide,
        "objetivos",
        slide_no,
        "La presentacion debe mostrar con claridad que problema se resuelve y con que objetivos",
    )
    _add_card(
        slide, 0.95, 1.48, 11.25, 1.32, fill=RGBColor(239, 246, 255), line=RGBColor(147, 197, 253)
    )
    _add_text(slide, "Objetivo general", 1.18, 1.74, 1.8, 0.2, size=14, color=BLUE, bold=True)
    _add_text(slide, GENERAL_OBJECTIVE, 1.18, 2.04, 10.6, 0.52, size=15, color=NAVY, bold=True)

    _add_card(slide, 0.95, 3.10, 5.45, 3.12, fill=WHITE, line=SLATE_LIGHT)
    _add_text(
        slide, "Objetivos especificos 1-3", 1.18, 3.38, 2.4, 0.2, size=14, color=BLUE, bold=True
    )
    _add_bullets(
        slide,
        [
            "Revisar sistematicamente calibracion, incertidumbre y CP en riesgo crediticio.",
            "Construir el dataset Lending Club 2007-2020 con variables objetivo PD, LGD y EAD.",
            "Disenar e implementar experimentos de CP sobre modelos base de PD, LGD y EAD.",
        ],
        1.18,
        3.74,
        4.95,
        2.10,
        size=13,
    )

    _add_card(slide, 6.75, 3.10, 5.45, 3.12, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(
        slide, "Objetivos especificos 4-6", 6.98, 3.38, 2.4, 0.2, size=14, color=TEAL, bold=True
    )
    _add_bullets(
        slide,
        [
            "Comparar CP con Platt, Isotonic y Bootstrap usando calibracion, cobertura y eficiencia.",
            "Evaluar impacto en IFRS 9, capital regulatorio y validacion interna.",
            "Proponer lineamientos de adopcion considerando implementacion, escalabilidad y gobernanza.",
        ],
        6.98,
        3.74,
        4.95,
        2.10,
        size=13,
    )
    _add_footer(
        slide,
        "Objetivos tomados de la propuesta y reorganizados visualmente para la presentacion oral.",
    )


def _slide_methodology_solution(slide, slide_no: int) -> None:
    _decorate_content_slide(
        slide,
        "metodologia",
        slide_no,
        "La metodologia de solucion propuesta conecta objetivos, datos, experimentos y adopcion",
    )
    steps = [
        (
            "Frente 1",
            "Estado del arte",
            ["Revisar literatura", "Comparar metodos", "Seleccionar tecnicas relevantes"],
            BLUE,
        ),
        (
            "Frente 2",
            "Datos",
            ["Construir dataset", "Definir PD/LGD/EAD", "Separar train-cal-test OOT"],
            TEAL,
        ),
        (
            "Frente 3",
            "Experimentacion",
            [
                "Entrenar y calibrar PD",
                "Aplicar Split / Mondrian / CQR",
                "Benchmark vs metodos clasicos",
            ],
            ORANGE,
        ),
        (
            "Frente 4",
            "Evaluacion y adopcion",
            ["Coverage y eficiencia", "IFRS9 y backtesting", "Recomendaciones de uso real"],
            GREEN,
        ),
    ]
    xs = [0.95, 3.85, 6.75, 9.65]
    for (tag, title, items, accent), x in zip(steps, xs, strict=False):
        _add_card(slide, x, 1.70, 2.40, 4.95, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
        _add_text(slide, tag, x + 0.18, 1.96, 0.8, 0.18, size=11, color=TEXT_SOFT, bold=True)
        _add_text(slide, title, x + 0.18, 2.22, 1.9, 0.28, size=16, color=accent, bold=True)
        current_y = 2.78
        for item in items:
            _add_card(slide, x + 0.14, current_y, 2.10, 0.70, fill=WHITE, line=SLATE_LIGHT)
            _add_text(
                slide,
                item,
                x + 0.26,
                current_y + 0.18,
                1.82,
                0.28,
                size=12,
                color=TEXT,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            current_y += 1.02

    for start_x in [3.35, 6.25, 9.15]:
        arrow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(start_x), Inches(3.74), Inches(0.26), Inches(0.42)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SLATE_LIGHT
        arrow.line.fill.background()

    _add_footer(
        slide,
        "Metodologia propuesta en cuatro frentes: revisar, construir, experimentar y evaluar/adoptar.",
    )


def _slide_context(slide, slide_no: int) -> None:
    _decorate_content_slide(
        slide,
        "introduccion",
        slide_no,
        "Introduccion: el proyecto no busca solo predecir mejor; busca decidir mejor bajo incertidumbre",
    )
    _add_metric(slide, 0.95, 1.46, 2.35, 1.0, "PD", "Probabilidad\nde default", accent=BLUE)
    _add_metric(slide, 3.50, 1.46, 2.35, 1.0, "LGD", "Severidad\nde perdida", accent=TEAL)
    _add_metric(
        slide, 6.05, 1.46, 2.35, 1.0, "EAD", "Exposicion al\nmomento del default", accent=ORANGE
    )
    _add_metric(slide, 8.60, 1.46, 2.55, 1.0, "Idea central", "ECL = PD x LGD x EAD", accent=GREEN)

    _add_card(
        slide, 0.95, 2.82, 6.25, 3.38, fill=RGBColor(239, 246, 255), line=RGBColor(147, 197, 253)
    )
    _add_text(
        slide, "Pregunta de investigacion", 1.18, 3.05, 2.2, 0.22, size=14, color=BLUE, bold=True
    )
    _add_text(slide, RESEARCH_QUESTION, 1.18, 3.38, 5.65, 1.18, size=16, color=NAVY, bold=True)
    _add_bullets(
        slide,
        [
            "El caso de estudio es Lending Club, un dataset publico y temporalmente rico.",
            "La historia integra modelado, incertidumbre y lectura de negocio en IFRS9.",
        ],
        1.18,
        4.85,
        5.5,
        1.0,
        size=14,
    )

    _add_card(slide, 7.45, 2.82, 4.72, 3.38, fill=WHITE, line=SLATE_LIGHT)
    _add_text(
        slide,
        "Por que esto importa para el grupo",
        7.7,
        3.05,
        3.0,
        0.22,
        size=14,
        color=TEAL,
        bold=True,
    )
    _add_bullets(
        slide,
        [
            "Riesgo de credito aparece en prestamos, tarjetas, provisiones y capital.",
            "Casi todos hemos visto scores; casi nadie ve la incertidumbre del score.",
            "Conformal Prediction permite contar esa incertidumbre con una garantia estadistica util.",
        ],
        7.7,
        3.40,
        4.0,
        2.2,
        size=15,
    )
    _add_footer(
        slide,
        "Contexto minimo para entrar al proyecto sin asumir experiencia previa en riesgo crediticio.",
    )


def _slide_conformal(slide, slide_no: int, chart: Path) -> None:
    _slide_state_of_art(slide, slide_no, chart)


def _slide_products(slide, slide_no: int) -> None:
    _decorate_content_slide(
        slide,
        "productos",
        slide_no,
        "El proyecto deja productos tecnicos y academicos, no solo resultados aislados",
    )
    products = [
        (
            "Producto tecnico",
            "Pagina Streamlit de tesis con narrativa, trazabilidad y soporte para sustentacion.",
            BLUE,
            0.95,
            1.62,
        ),
        (
            "Producto tecnico",
            "Pipeline reproducible con dataset temporal, benchmark PD/LGD/EAD e impacto IFRS9.",
            TEAL,
            6.25,
            1.62,
        ),
        (
            "Producto tecnico",
            "Presentacion estructurada para audiencia no experta con foco en problema y metodologia.",
            ORANGE,
            0.95,
            3.78,
        ),
        (
            "Producto academico",
            "Estado del arte, marco logico, OKR/KPI y matriz de riesgos construidos a partir de la propuesta.",
            GREEN,
            6.25,
            3.78,
        ),
    ]
    for tag, body, accent, x, y in products:
        _add_card(slide, x, y, 5.0, 1.72, fill=WHITE, line=SLATE_LIGHT)
        _add_text(slide, tag, x + 0.20, y + 0.20, 1.4, 0.18, size=11, color=accent, bold=True)
        _add_text(slide, body, x + 0.20, y + 0.48, 4.55, 0.82, size=14, color=NAVY, bold=True)

    _add_card(
        slide, 0.95, 5.92, 10.30, 0.58, fill=RGBColor(239, 246, 255), line=RGBColor(147, 197, 253)
    )
    _add_text(
        slide,
        "Lectura final: la tesis entrega metodo, evidencia, narrativa y artefactos utiles para continuar el trabajo.",
        1.20,
        6.10,
        9.8,
        0.20,
        size=12,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_footer(
        slide, "Se incluye una seccion explicita de productos, tal como sugirio el feedback."
    )


def _slide_conformal_original(slide, slide_no: int, chart: Path) -> None:
    _decorate_content_slide(slide, "concepto", slide_no, "Conformal Prediction en una imagen")
    _add_card(slide, 0.95, 1.44, 7.1, 4.98, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 1.15, 1.76, 6.65, 4.35)

    _add_card(slide, 8.32, 1.44, 3.9, 4.98, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Lectura guiada", 8.56, 1.72, 1.8, 0.2, size=14, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "Primero entrenamos y calibramos el modelo puntual.",
            "Luego usamos un set de calibracion para medir errores pasados.",
            "Ese error se resume en un cuantil q al nivel objetivo (90% o 95%).",
            "Finalmente, cada nueva prediccion recibe un intervalo plausible en vez de un solo numero.",
            "En Mondrian, ese q se calcula por grade para no esconder subgrupos mal cubiertos.",
        ],
        8.56,
        2.08,
        3.4,
        2.45,
        size=13,
    )
    _add_text(
        slide,
        "Calibracion y conformal no compiten",
        8.56,
        5.14,
        2.8,
        0.2,
        size=13,
        color=TEAL,
        bold=True,
    )
    _add_text(
        slide,
        "La calibracion corrige el nivel probabilistico; conformal agrega una banda de incertidumbre con cobertura controlada.",
        8.56,
        5.42,
        3.35,
        0.62,
        size=12,
        color=TEXT_SOFT,
    )
    _add_footer(
        slide,
        "Figura propia inspirada en Angelopoulos y Bates (2023) y en el marco formal de ALRW (2022).",
    )


def _slide_dataset(slide, slide_no: int, splits_chart: Path, grade_chart: Path, data: dict) -> None:
    _decorate_content_slide(
        slide,
        "metodologia",
        slide_no,
        "Diseno experimental: dataset, horizonte temporal y señal de riesgo del caso de estudio",
    )
    shapes = data["shapes"]
    eda = data["eda"]
    split_info = data["split_info"]
    cleaned_rows = int(shapes["data/interim/lending_club_cleaned.parquet"]["rows"])
    n_cols = int(shapes["data/interim/lending_club_cleaned.parquet"]["cols"])

    _add_metric(slide, 0.95, 1.48, 1.8, 0.96, "Limpio", _format_count(cleaned_rows), accent=BLUE)
    _add_metric(slide, 0.95, 2.60, 1.8, 0.96, "Columnas", str(n_cols), accent=TEAL)
    _add_metric(slide, 0.95, 3.72, 1.8, 0.96, "Horizonte", "2007-2020", accent=ORANGE)

    _add_card(slide, 2.95, 1.48, 4.05, 2.15, fill=WHITE, line=SLATE_LIGHT)
    _add_text(
        slide,
        "Que dataset usamos exactamente",
        3.18,
        1.72,
        2.5,
        0.2,
        size=14,
        color=BLUE,
        bold=True,
    )
    _add_bullets(
        slide,
        [
            "Prestamos historicos de Lending Club, ya curados para analitica.",
            f"Se limpian a {_format_count(cleaned_rows)} prestamos y {n_cols} columnas.",
            "Luego se separan train, calibration y test de manera temporal para evitar leakage.",
        ],
        3.18,
        2.03,
        3.45,
        1.3,
        size=13,
    )

    _add_card(slide, 2.95, 3.86, 4.05, 2.42, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(
        slide,
        "Por que si sirve para esta tesis",
        3.18,
        4.10,
        2.3,
        0.2,
        size=14,
        color=TEAL,
        bold=True,
    )
    _add_bullets(
        slide,
        [
            f"El default rate escala de A={eda['default_rate_by_grade']['A'] * 100:.1f}% a G={eda['default_rate_by_grade']['G'] * 100:.1f}%.",
            "Ese gradiente confirma senal de negocio, no solo volumen.",
            "Al ser publico, el proyecto queda replicable y defendible para aula y jurados.",
        ],
        3.18,
        4.40,
        3.45,
        1.45,
        size=13,
    )

    _add_card(slide, 7.28, 1.48, 4.85, 2.2, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, splits_chart, 7.48, 1.70, 4.45, 1.78)
    _add_card(slide, 7.28, 3.86, 4.85, 2.42, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, grade_chart, 7.48, 4.06, 4.45, 2.00)

    split_text = " | ".join(f"{row['label']}={_format_count(row['n'])}" for row in split_info)
    _add_footer(slide, f"Diseno experimental con splits OOT del proyecto: {split_text}.")


def _slide_pipeline(slide, slide_no: int, pipeline_chart: Path) -> None:
    _decorate_content_slide(
        slide,
        "metodologia",
        slide_no,
        "Pipeline de ejecucion: de la propuesta metodologica a una corrida trazable de punta a punta",
    )
    _add_card(slide, 0.95, 1.42, 11.25, 4.65, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, pipeline_chart, 1.15, 1.74, 10.85, 3.55)

    _add_card(
        slide, 0.95, 6.15, 3.45, 0.72, fill=RGBColor(239, 246, 255), line=RGBColor(147, 197, 253)
    )
    _add_text(
        slide,
        "1. Separar train / calibration / test no es detalle tecnico; evita leakage.",
        1.16,
        6.32,
        3.0,
        0.25,
        size=11,
        color=NAVY,
        bold=True,
    )

    _add_card(
        slide, 4.55, 6.15, 3.45, 0.72, fill=RGBColor(240, 253, 250), line=RGBColor(94, 234, 212)
    )
    _add_text(
        slide,
        "2. La capa conformal se monta sobre PD calibrada; no reemplaza calibracion.",
        4.76,
        6.32,
        3.0,
        0.25,
        size=11,
        color=NAVY,
        bold=True,
    )

    _add_card(
        slide, 8.15, 6.15, 4.05, 0.72, fill=RGBColor(255, 247, 237), line=RGBColor(253, 186, 116)
    )
    _add_text(
        slide,
        "3. El cierre no se queda en ML: termina en backtesting y lectura de negocio.",
        8.36,
        6.32,
        3.55,
        0.25,
        size=11,
        color=NAVY,
        bold=True,
    )

    _add_footer(
        slide,
        "Pipeline de especializacion: alcance acotado, ejecutable y alineado con una lectura CRISP-DM del trabajo.",
    )


def _slide_crisp(slide, slide_no: int) -> None:
    _decorate_content_slide(
        slide,
        "metodologia",
        slide_no,
        "CRISP-DM se uso como marco de ejecucion, no solo como formalidad documental",
    )
    accent_cycle = [BLUE, TEAL, ORANGE, BLUE, TEAL, ORANGE]
    positions = [
        (0.95, 1.52),
        (4.17, 1.52),
        (7.39, 1.52),
        (0.95, 4.02),
        (4.17, 4.02),
        (7.39, 4.02),
    ]
    for (x, y), phase, accent in zip(positions, CRISP_PHASES, accent_cycle, strict=False):
        _add_phase_card(slide, x, y, 2.86, 2.05, phase["title"], phase["body"], accent)

    _add_card(slide, 10.45, 1.52, 1.73, 4.55, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Lectura extendida", 10.68, 1.82, 1.2, 0.2, size=13, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "Cada fase deja una evidencia concreta.",
            "Eso evita historias bonitas sin soporte tecnico.",
            "Tambien ayuda a defender el alcance real del proyecto.",
        ],
        10.68,
        2.20,
        1.2,
        2.2,
        size=12,
    )
    _add_text(
        slide,
        "Resultado global: una tesis trazable de punta a punta, no una coleccion de experimentos aislados.",
        10.68,
        4.95,
        1.2,
        0.7,
        size=11,
        color=TEXT_SOFT,
        bold=True,
    )
    _add_footer(slide, "Lectura extendida de la tabla CRISP-DM mostrada en la pagina de tesis.")


def _slide_objectives(slide, slide_no: int, blocks: list[dict], title: str, footer: str) -> None:
    _decorate_content_slide(slide, "trazabilidad", slide_no, title)
    accents = [BLUE, TEAL, ORANGE]
    y_positions = [1.48, 3.20, 4.92]
    for y, block, accent in zip(y_positions, blocks, accents, strict=False):
        _add_objective_card(
            slide,
            0.95,
            y,
            11.25,
            1.46,
            block["objective"],
            block["deliverable"],
            block["result"],
            accent,
        )
    _add_footer(slide, footer)


def _slide_pd(slide, slide_no: int, chart: Path, model: dict) -> None:
    _decorate_content_slide(
        slide,
        "resultados",
        slide_no,
        "PD: la calibracion mejora la utilidad probabilistica sin perder discriminacion",
    )
    final_model = model["models"][3]
    _add_metric(slide, 0.95, 1.48, 1.7, 0.96, "AUC OOT", f"{final_model['auc']:.3f}", accent=BLUE)
    _add_metric(slide, 0.95, 2.58, 1.7, 0.96, "Brier", f"{final_model['brier']:.3f}", accent=TEAL)
    _add_metric(slide, 0.95, 3.68, 1.7, 0.96, "Calibrador", "Isotonic", accent=ORANGE)

    _add_card(slide, 2.90, 1.48, 6.25, 4.82, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 3.12, 1.85, 5.8, 4.15)

    _add_card(slide, 9.40, 1.48, 2.80, 4.82, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Lectura", 9.62, 1.80, 1.1, 0.2, size=14, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "El baseline logistico sirve como piso interpretable.",
            "CatBoost calibrado mantiene AUC y mejora calidad probabilistica.",
            "Eso importa porque en riesgo una PD mal calibrada se convierte en una mala decision.",
        ],
        9.62,
        2.16,
        2.30,
        2.25,
        size=14,
    )
    _add_text(slide, "Mensaje simple", 9.62, 5.02, 1.5, 0.2, size=13, color=TEAL, bold=True)
    _add_text(
        slide,
        "Pasamos de score util para ranking a probabilidad mas util para originacion, pricing y provisiones.",
        9.62,
        5.28,
        2.30,
        0.55,
        size=12,
        color=TEXT_SOFT,
    )
    _add_footer(slide, "Comparativa OOT de PD: baseline logistico vs CatBoost calibrado.")


def _slide_coverage_target(
    slide, slide_no: int, chart: Path, conformal: dict, monthly: pd.DataFrame
) -> None:
    _decorate_content_slide(
        slide,
        "resultados",
        slide_no,
        "Cobertura objetivo: que significa y como se verifico en el tiempo",
    )
    _add_card(slide, 0.95, 1.48, 7.15, 4.78, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 1.14, 1.82, 6.75, 4.05)

    _add_metric(
        slide,
        8.35,
        1.48,
        1.65,
        0.96,
        "Coverage 90",
        f"{conformal['coverage_90'] * 100:.2f}%",
        accent=BLUE,
    )
    _add_metric(
        slide,
        10.12,
        1.48,
        1.65,
        0.96,
        "Coverage 95",
        f"{conformal['coverage_95'] * 100:.2f}%",
        accent=TEAL,
    )
    _add_metric(slide, 8.35, 2.58, 1.65, 0.96, "Meses", str(len(monthly)), accent=ORANGE)
    _add_metric(
        slide,
        10.12,
        2.58,
        1.65,
        0.96,
        "Checks",
        f"{conformal['checks_passed']}/{conformal['checks_total']}",
        accent=GREEN,
    )

    _add_card(slide, 8.35, 3.78, 3.42, 2.48, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Como leer la cobertura", 8.58, 4.04, 2.3, 0.2, size=14, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "Meta 90%: cerca de 9 de cada 10 observaciones nuevas deben caer dentro del intervalo.",
            "Es una garantia marginal en el set OOT, no una promesa perfecta para cada grade o cada mes.",
            "Por eso, ademas del promedio global, medimos cobertura por grade y backtest temporal.",
        ],
        8.58,
        4.36,
        2.80,
        1.55,
        size=12,
    )
    _add_footer(
        slide,
        "Backtest temporal PD: la cobertura global cumple, aunque siguen advertencias estadisticas abiertas para monitoreo.",
    )


def _slide_group_coverage(slide, slide_no: int, chart: Path, group_metrics: pd.DataFrame) -> None:
    _decorate_content_slide(
        slide,
        "resultados",
        slide_no,
        "Cobertura por grade: donde la incertidumbre realmente cambia",
    )
    df = group_metrics.rename(columns={"group": "grade"}).copy()
    worst = df.loc[df["coverage_90"].idxmin()]
    best = df.loc[df["coverage_90"].idxmax()]
    under_target = int((df["coverage_90"] < 0.90).sum())

    _add_card(slide, 0.95, 1.48, 7.15, 4.82, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 1.14, 1.84, 6.75, 4.10)

    _add_metric(
        slide,
        8.35,
        1.48,
        1.65,
        0.96,
        "Peor grade",
        f"{worst['grade']}  {worst['coverage_90'] * 100:.1f}%",
        accent=ORANGE,
    )
    _add_metric(
        slide,
        10.12,
        1.48,
        1.65,
        0.96,
        "Mejor grade",
        f"{best['grade']}  {best['coverage_90'] * 100:.1f}%",
        accent=BLUE,
    )
    _add_metric(slide, 8.35, 2.58, 1.65, 0.96, "Bajo meta", str(under_target), accent=RED)
    _add_metric(slide, 10.12, 2.58, 1.65, 0.96, "Grades", str(len(df)), accent=TEAL)

    _add_card(slide, 8.35, 3.78, 3.42, 2.52, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Por que medirlo asi", 8.58, 4.05, 2.0, 0.2, size=14, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "El portafolio esta desbalanceado: B y C pesan mucho mas que F y G.",
            "Si solo miro la cobertura promedio, puedo ocultar fallos en segmentos de alto riesgo.",
            "Coverage por grade muestra si la garantia sirve donde mas duele equivocarse.",
        ],
        8.58,
        4.36,
        2.80,
        1.55,
        size=12,
    )
    _add_footer(
        slide,
        "Mondrian deja visible la heterogeneidad por segmento; esa es la antesala del contraste con el enfoque global.",
    )


def _slide_global_vs_mondrian(slide, slide_no: int, chart: Path, benchmark: pd.DataFrame) -> None:
    _decorate_content_slide(
        slide,
        "resultados",
        slide_no,
        "Global split vs Mondrian: diferencia y utilidad en este proyecto",
    )
    bench = benchmark.set_index("variant")
    global_row = bench.loc["global_split"]
    mondrian_row = bench.loc["mondrian_selected_cfg"]

    _add_card(slide, 0.95, 1.48, 7.05, 4.82, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 1.14, 1.84, 6.65, 4.08)

    _add_card(slide, 8.25, 1.48, 3.95, 2.08, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Tradicional / global", 8.48, 1.76, 1.6, 0.2, size=14, color=ORANGE, bold=True)
    _add_bullets(
        slide,
        [
            "Sirve como baseline simple y rapido.",
            "Da una garantia promedio del portafolio.",
            "Problema: puede dejar grades pequenos o riesgosos muy mal cubiertos.",
        ],
        8.48,
        2.06,
        3.15,
        1.18,
        size=12,
    )

    _add_card(
        slide, 8.25, 3.78, 3.95, 2.52, fill=RGBColor(240, 253, 250), line=RGBColor(94, 234, 212)
    )
    _add_text(slide, "Mondrian", 8.48, 4.05, 1.0, 0.2, size=14, color=TEAL, bold=True)
    _add_bullets(
        slide,
        [
            "Calcula el cuantil por grade, no solo una vez para todo el portafolio.",
            f"Min group coverage sube de {global_row['min_group_coverage'] * 100:.1f}% a {mondrian_row['min_group_coverage'] * 100:.1f}%.",
            f"Ademas, el ancho promedio baja de {global_row['avg_width']:.3f} a {mondrian_row['avg_width']:.3f}.",
        ],
        8.48,
        4.36,
        3.15,
        1.45,
        size=12,
    )
    _add_footer(
        slide,
        "Utilidad practica: global para benchmark y sanity check; Mondrian para una garantia operativa por segmento en un portafolio desbalanceado.",
    )


def _slide_lgd_ead(slide, slide_no: int, chart: Path, status: dict) -> None:
    _decorate_content_slide(
        slide,
        "resultados",
        slide_no,
        "LGD y EAD: la triada PD-LGD-EAD ya tiene evidencia ejecutable",
    )
    lgd = status["lgd"]
    ead = status["ead"]
    _add_metric(
        slide, 0.95, 1.48, 1.75, 0.96, "Train", _format_count(int(lgd["n_train"])), accent=BLUE
    )
    _add_metric(
        slide, 0.95, 2.58, 1.75, 0.96, "Calibration", _format_count(int(lgd["n_cal"])), accent=TEAL
    )
    _add_metric(
        slide, 0.95, 3.68, 1.75, 0.96, "Test", _format_count(int(lgd["n_test"])), accent=ORANGE
    )

    _add_card(slide, 2.95, 1.48, 6.15, 4.82, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 3.15, 1.84, 5.75, 4.08)

    _add_card(slide, 9.35, 1.48, 2.85, 4.82, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Lectura", 9.58, 1.78, 1.0, 0.2, size=14, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            f"Se probaron 4 variantes LGD; la seleccionada fue {lgd['selected_variant']}.",
            "Solo la variante adaptive grade-temporal paso coverage, min grade y sesgo.",
            f"EAD tambien cumple: {ead['conformal']['metrics_90']['empirical_coverage'] * 100:.2f}% al 90% y {ead['conformal']['metrics_95']['empirical_coverage'] * 100:.2f}% al 95%.",
        ],
        9.58,
        2.12,
        2.35,
        2.55,
        size=12,
    )
    _add_text(
        slide,
        "Mensaje simple: el proyecto ya no habla solo de PD; puede contar la triada completa con evidencia real.",
        9.58,
        5.22,
        2.35,
        0.62,
        size=12,
        color=TEXT_SOFT,
        bold=True,
    )
    _add_footer(
        slide,
        "LGD/EAD sobre subset defaults-only, con benchmark de variantes y seleccion por guardrails.",
    )


def _slide_ifrs9(slide, slide_no: int, chart: Path, ifrs9: pd.DataFrame) -> None:
    _decorate_content_slide(
        slide, "impacto", slide_no, "IFRS9: la incertidumbre si cambia la lectura de provisiones"
    )
    baseline = ifrs9.loc[ifrs9["scenario"] == "baseline"].iloc[0]
    severe = ifrs9.loc[ifrs9["scenario"] == "severe"].iloc[0]
    uplift = (float(severe["total_ecl"]) / float(baseline["total_ecl"]) - 1.0) * 100

    _add_metric(
        slide, 0.95, 1.48, 1.9, 0.96, "Baseline", f"{baseline['total_ecl'] / 1e9:.2f}B", accent=BLUE
    )
    _add_metric(
        slide, 0.95, 2.58, 1.9, 0.96, "Severe", f"{severe['total_ecl'] / 1e9:.2f}B", accent=ORANGE
    )
    _add_metric(slide, 0.95, 3.68, 1.9, 0.96, "Uplift", f"+{uplift:.0f}%", accent=TEAL)

    _add_card(slide, 3.10, 1.48, 5.75, 4.82, fill=WHITE, line=SLATE_LIGHT)
    _add_picture(slide, chart, 3.35, 1.88, 5.25, 4.00)

    _add_card(slide, 9.15, 1.48, 3.05, 4.82, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Lo importante aqui", 9.38, 1.78, 1.7, 0.2, size=14, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "La tesis no se queda en metricas de machine learning.",
            "La incertidumbre llega hasta una lectura financiera concreta: provisiones.",
            "Ese salto es clave para hablarle a negocio y no solo a cientificos de datos.",
        ],
        9.38,
        2.12,
        2.55,
        1.85,
        size=13,
    )
    _add_text(slide, "Stages baseline", 9.38, 4.66, 1.2, 0.2, size=13, color=TEAL, bold=True)
    _add_text(
        slide,
        "S1 34.1% | S2 43.4% | S3 22.5%",
        9.38,
        4.94,
        2.2,
        0.25,
        size=12,
        color=TEXT_SOFT,
        bold=True,
    )
    _add_footer(
        slide,
        "Impacto IFRS9 del pipeline canonico: el ECL reacciona de forma material al estres macroeconomico.",
    )


def _slide_conclusion(slide, slide_no: int, fairness: dict, governance: dict) -> None:
    _decorate_content_slide(slide, "cierre", slide_no, "Que quedo demostrado y que sigue abierto")
    _add_card(slide, 0.95, 1.52, 5.65, 4.86, fill=WHITE, line=SLATE_LIGHT)
    _add_text(slide, "Logros del proyecto", 1.18, 1.84, 2.0, 0.2, size=15, color=BLUE, bold=True)
    _add_bullets(
        slide,
        [
            "La tesis explica riesgo crediticio a una audiencia no experta sin perder rigor.",
            "PD queda calibrada y con incertidumbre interpretable.",
            "LGD y EAD ya alcanzan cobertura defendible con benchmark explicito.",
            "IFRS9 conecta la parte tecnica con una lectura economica real.",
        ],
        1.18,
        2.18,
        5.0,
        2.45,
        size=15,
    )

    _add_card(slide, 6.90, 1.52, 5.30, 4.86, fill=RGBColor(248, 250, 252), line=SLATE_LIGHT)
    _add_text(slide, "Temas abiertos", 7.14, 1.84, 1.8, 0.2, size=15, color=ORANGE, bold=True)
    gov = "PASS" if governance["overall_pass"] else "WARN"
    _add_bullets(
        slide,
        [
            f"Fairness pasa en {fairness['n_passed']}/{fairness['n_attributes']} atributos evaluados.",
            f"Governance sigue en {gov}; todavia hay agenda de monitoreo y drift.",
            "La continuidad natural es robustecer la politica operativa y escalar experimentos avanzados.",
        ],
        7.14,
        2.18,
        4.7,
        2.05,
        size=15,
    )
    _add_text(
        slide,
        "Mensaje final: Conformal Prediction sirve porque hace visible lo que antes el modelo callaba.",
        7.14,
        5.04,
        4.7,
        0.62,
        size=15,
        color=NAVY,
        bold=True,
    )
    _add_footer(slide, "Slide de cierre para dejar visible durante preguntas.")


def _slide_references(slide, slide_no: int) -> None:
    _decorate_content_slide(slide, "referencias", slide_no, "Referencias clave del proyecto")
    _add_card(slide, 0.95, 1.48, 12.0, 5.28, fill=WHITE, line=SLATE_LIGHT)
    _add_bullets(slide, REFERENCES, 1.22, 1.86, 11.45, 4.8, size=15, color=TEXT)
    _add_footer(slide, "Base teorica y normativa usada para la narrativa del proyecto.")


def build_deck(output_path: Path, assets_dir: Path) -> Path:
    root = Path(__file__).resolve().parents[1]
    charts, data = _generate_assets(root, assets_dir)

    model = _load_json(root / "data/processed/model_comparison.json")
    conformal = _load_json(root / "models/conformal_policy_status.json")
    fairness = _load_json(root / "models/fairness_audit_status.json")
    governance = _load_json(root / "models/governance_status.json")

    prs = _new_presentation()
    blank = prs.slide_layouts[6]

    slides = [
        lambda slide, _: _title_slide(slide),
        _agenda_slide,
        _slide_context,
        _slide_justification,
        lambda slide, no: _slide_conformal(slide, no, charts["conformal"]),
        _slide_objectives_overview,
        _slide_methodology_solution,
        lambda slide, no: _slide_dataset(slide, no, charts["splits"], charts["grade"], data),
        lambda slide, no: _slide_pipeline(slide, no, charts["pipeline"]),
        lambda slide, no: _slide_pd(slide, no, charts["pd"], model),
        lambda slide, no: _slide_coverage_target(
            slide, no, charts["coverage"], conformal, data["monthly"]
        ),
        lambda slide, no: _slide_global_vs_mondrian(
            slide, no, charts["mondrian"], data["benchmark"]
        ),
        lambda slide, no: _slide_lgd_ead(slide, no, charts["lgd_ead"], data["lgd_ead"]),
        lambda slide, no: _slide_ifrs9(slide, no, charts["ifrs9"], data["ifrs9"]),
        _slide_products,
        lambda slide, no: _slide_conclusion(slide, no, fairness, governance),
        _slide_references,
    ]

    for idx, builder in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        builder(slide, idx)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> None:
    args = _parse_args()
    out = build_deck(Path(args.output), Path(args.assets_dir))
    print(out)


if __name__ == "__main__":
    main()
