"""Build improved thesis defense PPTX for Especialización."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── Colors ──────────────────────────────────────────────
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)      # Dark navy
ACCENT = RGBColor(0x2E, 0x75, 0xB6)       # Mid-blue
BODY = RGBColor(0x2D, 0x2D, 0x2D)         # Near-black
MUTED = RGBColor(0x77, 0x77, 0x77)        # Gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xEB, 0xF3, 0xFA)     # Light blue bg
RULE = RGBColor(0xCC, 0xCC, 0xCC)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF2, 0xCC)  # Yellow callout
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
DARK_BG = RGBColor(0x1F, 0x4E, 0x79)
SECTION_BG = RGBColor(0x1A, 0x3A, 0x5C)

FONT_FACE = "Arial"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── Paths ───────────────────────────────────────────────
ROOT = Path("/home/eigenlinux/projects/lending-club-risk-project")
FIG_DIR = ROOT / "reports" / "paper_material" / "figures_publication"
NB_IMG = ROOT / "dist" / "streamlit_deploy" / "reports" / "notebook_images"
OUT = ROOT / "thesis_poster" / "Presentacion_Tesis_Carlos_Final.pptx"


def _add_text(slide, text, x, y, w, h, *, size=20, bold=False, color=BODY,
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT_FACE):
    """Helper to add a text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    tf.auto_size = None
    return txBox


def _add_multiline(slide, lines, x, y, w, h, *, size=18, color=BODY, bullet=False,
                   spacing=Pt(6), bold_prefix=False):
    """Add multiple lines with optional bold prefix (text before ':')."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.level = 0
        if bold_prefix and ":" in line:
            prefix, rest = line.split(":", 1)
            run1 = p.add_run()
            run1.text = prefix + ":"
            run1.font.bold = True
            run1.font.size = Pt(size)
            run1.font.color.rgb = color
            run1.font.name = FONT_FACE
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.name = FONT_FACE
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = FONT_FACE
    return txBox


def _add_rule(slide, x, y, w, h=0.03, color=RULE):
    """Add a thin horizontal rule."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _dark_slide(slide):
    """Set dark background."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG


def _section_slide(slide, label, title):
    """Create a section divider slide."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = SECTION_BG
    _add_text(slide, label, 0.5, 1.8, 9.0, 0.4, size=16, color=RGBColor(0x7B, 0xAF, 0xD4))
    _add_text(slide, title, 0.5, 2.2, 9.0, 1.0, size=36, bold=True, color=WHITE)
    _add_rule(slide, 0.5, 3.3, 2.5, 0.06, ACCENT)


def _kpi_box(slide, x, y, label, value, *, w=2.0, h=0.9, color=ACCENT):
    """Add a KPI metric box."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.5)
    # Label
    _add_text(slide, label, x + 0.1, y + 0.05, w - 0.2, 0.3, size=11, color=MUTED,
              align=PP_ALIGN.CENTER)
    # Value
    _add_text(slide, value, x + 0.1, y + 0.35, w - 0.2, 0.45, size=22, bold=True,
              color=color, align=PP_ALIGN.CENTER)


def _add_image_safe(slide, path, x, y, w, h=None):
    """Add image if file exists."""
    p = Path(path)
    if p.exists():
        kwargs = {"left": Inches(x), "top": Inches(y), "width": Inches(w)}
        if h:
            kwargs["height"] = Inches(h)
        slide.shapes.add_picture(str(p), **kwargs)
        return True
    return False


def _add_table(slide, headers, rows, x, y, w, row_h=0.35, *, header_color=PRIMARY,
               font_size=12):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_w = w / n_cols
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                          Inches(w), Inches(row_h * n_rows))
    table = table_shape.table
    # Headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FACE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
    # Data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = BODY
                p.font.name = FONT_FACE
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    return table_shape


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank layout

    # ════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _add_text(s, "Conformal Prediction para la Calibracion\ny Cuantificacion de Incertidumbre en\nModelos de Riesgo Crediticio",
              0.7, 0.6, 8.6, 2.0, size=30, bold=True, color=WHITE)
    _add_rule(s, 0.7, 2.7, 2.5, 0.05, ACCENT)
    _add_text(s, "Tesis de Especializacion en Analitica y Ciencia de Datos Aplicada",
              0.7, 2.85, 8.6, 0.4, size=16, color=RGBColor(0xA0, 0xBB, 0xDD))
    _add_text(s, "Carlos Alfredo Vergara Rojas",
              0.7, 3.4, 8.6, 0.35, size=18, color=RGBColor(0xCA, 0xDC, 0xFC))
    _add_text(s, "Directora: Eliana Mirledy Toro Ocampo",
              0.7, 3.8, 8.6, 0.35, size=15, color=RGBColor(0xCA, 0xDC, 0xFC))
    _add_text(s, "Universidad Tecnologica de Pereira  |  Marzo 2026",
              0.7, 4.3, 8.6, 0.35, size=14, color=RGBColor(0xA0, 0xBB, 0xDD))
    _add_text(s, "\"Una probabilidad puntual ordena clientes; una prediccion con incertidumbre permite decidir con prudencia.\"",
              0.7, 4.85, 8.6, 0.5, size=13, color=RGBColor(0x88, 0xAA, 0xCC))

    # ════════════════════════════════════════════════════════
    # SLIDE 2: AGENDA
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Estructura de la presentacion", 0.5, 0.2, 9.0, 0.7, size=26, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 0.85, 9.0)
    items = [
        "1. Motivacion y planteamiento del problema",
        "2. Pregunta de investigacion y objetivos",
        "3. Marco teorico: Conformal Prediction",
        "4. Metodologia (6 fases CRISP-DM)",
        "5. Resultados: PD, calibracion, Mondrian, LGD/EAD",
        "6. Impacto regulatorio IFRS 9",
        "7. Respuesta a la evaluacion del anteproyecto",
        "8. Cumplimiento de objetivos y conclusiones",
    ]
    _add_multiline(s, items, 0.7, 1.0, 8.5, 4.0, size=20)

    # ════════════════════════════════════════════════════════
    # SLIDE 3: MOTIVATION - ECL formula
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Las decisiones de credito dependen de PD, LGD y EAD,\npero hoy se toman con incertidumbre oculta",
              0.5, 0.2, 9.0, 0.8, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # KPIs
    _kpi_box(s, 0.7, 1.2, "PD", "Prob. default")
    _kpi_box(s, 3.0, 1.2, "LGD", "Severidad")
    _kpi_box(s, 5.3, 1.2, "EAD", "Exposicion")
    _add_text(s, "ECL = PD x LGD x EAD", 7.6, 1.35, 2.2, 0.6, size=20, bold=True,
              color=PRIMARY, align=PP_ALIGN.CENTER)

    _add_multiline(s, [
        "Estas metricas alimentan provisiones (IFRS 9) y capital regulatorio (Basilea III).",
        "Los modelos ML logran buen ranking (AUC-ROC), pero producen probabilidades mal calibradas.",
        "Una PD de 12% no indica si el valor real esta entre 10-14% o entre 5-28%.",
        "Sin incertidumbre explicita: perdidas inesperadas, provisiones excesivas, desconfianza regulatoria.",
    ], 0.7, 2.3, 8.5, 2.8, size=18, bold_prefix=False)

    _add_text(s, "Fuente: Basel Committee (2006), IASB (2014), Angelopoulos & Bates (2023)",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 4: RESEARCH QUESTION
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Pregunta de investigacion y contribucion central del proyecto",
              0.5, 0.2, 9.0, 0.7, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 0.85, 9.0)

    # Research question box
    from pptx.enum.shapes import MSO_SHAPE
    rq_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.1),
                                 Inches(8.0), Inches(1.4))
    rq_box.fill.solid()
    rq_box.fill.fore_color.rgb = LIGHT_BG
    rq_box.line.color.rgb = ACCENT
    rq_box.line.width = Pt(2)
    _add_text(s, "Como mejorar la calibracion y la cuantificacion de la incertidumbre\nen los modelos de riesgo crediticio (PD, LGD, EAD) mediante la aplicacion\nde tecnicas de prediccion conformal, y cual es su desempeno\ncomparado frente a metodos tradicionales de calibracion?",
              1.2, 1.2, 7.6, 1.2, size=17, color=PRIMARY, align=PP_ALIGN.CENTER)

    _add_text(s, "Contribucion", 0.5, 2.7, 9.0, 0.35, size=20, bold=True, color=ACCENT)
    _add_multiline(s, [
        "Extension de CP a la triada completa PD-LGD-EAD (literatura solo cubre PD).",
        "Demostracion de +30pp en cobertura minima con Mondrian vs Split Global.",
        "Propuesta del ancho conformal como senal SICR para clasificacion por etapas IFRS 9.",
    ], 0.7, 3.1, 8.5, 1.8, size=17, bold_prefix=False)

    _add_text(s, "Caso de estudio: Lending Club (2007-2020), 1.86M prestamos, dataset publico y replicable.",
              0.5, 4.85, 9.0, 0.5, size=14, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 5: OBJECTIVES (reformulated from document)
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Seis objetivos especificos con alcance acotado, medible y cumplido",
              0.5, 0.2, 9.0, 0.7, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 0.85, 9.0)

    _add_table(s, ["OE", "Alcance planificado", "Resultado obtenido"],
               [
                   ["OE1", "Marco teorico CP (revision sistematica)", "19 fuentes, vacio PD-only identificado"],
                   ["OE2", "Dataset OOT sin leakage (PD/LGD/EAD)", "1.86M prestamos, 42 features, 3 splits"],
                   ["OE3", "PD calibrado (AUC, Brier, ECE)", "AUC 0.713, Brier 0.155, ECE 0.007"],
                   ["OE4", "CP Mondrian PD + variantes LGD/EAD", "PD 92.42%, LGD 90.50%, EAD 91.20%"],
                   ["OE5", "Impacto IFRS 9 (ECL, stages, SICR)", "$1,001M-$1,799M, SICR conformal"],
                   ["OE6", "Lineamientos de adopcion practica", "Backtesting 35m, gate 7/13, dashboard"],
               ], 0.3, 0.95, 9.4, row_h=0.42, font_size=12)

    _add_text(s, "Todos los objetivos cumplidos con evidencia cuantitativa verificable (Tabla 12 del documento tecnico).",
              0.5, 4.1, 9.0, 0.35, size=14, bold=True, color=GREEN)

    _add_text(s, "Cambio vs. anteproyecto: se acoto alcance a CP sobre PD/LGD/EAD e IFRS 9, excluyendo optimizacion robusta e inferencia causal (reservados para maestria).",
              0.5, 4.55, 9.0, 0.5, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 6: WHY CP - conceptual
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Conformal Prediction ofrece garantias formales de cobertura\nque bootstrap y metodos bayesianos no proveen",
              0.5, 0.2, 9.0, 0.8, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # Two columns
    _add_text(s, "Sin incertidumbre", 0.6, 1.2, 4.0, 0.35, size=18, bold=True, color=RED)
    _add_multiline(s, [
        "PD = 12%",
        "El modelo no muestra si esa cifra es estable o fragil.",
        "Sobreconfianza en originacion y provision.",
    ], 0.6, 1.6, 4.0, 1.5, size=16)

    _add_text(s, "Con Conformal Prediction", 5.3, 1.2, 4.5, 0.35, size=18, bold=True, color=GREEN)
    _add_multiline(s, [
        "PD = 12%, intervalo [8%, 17%]",
        "Garantia: P(Y in C(X)) >= 1-alpha",
        "Sin supuestos parametricos (distribution-free).",
        "Solo requiere intercambiabilidad de datos.",
    ], 5.3, 1.6, 4.5, 1.8, size=16)

    _add_text(s, "Ventajas de CP vs. alternativas", 0.6, 3.4, 9.0, 0.35, size=18, bold=True, color=ACCENT)
    _add_multiline(s, [
        "vs. Bootstrap: CP tiene garantias en muestras finitas; bootstrap no.",
        "vs. Bayesiano: CP no requiere distribuciones a priori.",
        "vs. Calibracion sola: la calibracion corrige nivel probabilistico; CP agrega banda de incertidumbre.",
    ], 0.6, 3.8, 8.8, 1.5, size=16, bold_prefix=True)

    _add_text(s, "Vovk et al. (2022); Angelopoulos & Bates (2023); Romano et al. (2019)",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 7: CP VISUAL EXAMPLE
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Cada prediccion se acompana de un intervalo cuyo ancho\nrefleja la confianza del modelo en esa estimacion",
              0.5, 0.2, 9.0, 0.8, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # Three client examples
    examples = [
        ("Cliente A", "PD 6%", "[4%, 8%]", "Baja incertidumbre", GREEN),
        ("Cliente B", "PD 12%", "[8%, 17%]", "Media incertidumbre", ACCENT),
        ("Cliente C", "PD 14%", "[5%, 28%]", "Alta incertidumbre", RED),
    ]
    for i, (name, pd_val, interval, label, col) in enumerate(examples):
        x = 0.6 + i * 3.1
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2),
                                  Inches(2.8), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = col
        box.line.width = Pt(2)
        _add_text(s, name, x + 0.1, 1.3, 2.6, 0.3, size=16, bold=True, color=PRIMARY,
                  align=PP_ALIGN.CENTER)
        _add_text(s, pd_val, x + 0.1, 1.65, 2.6, 0.35, size=22, bold=True, color=col,
                  align=PP_ALIGN.CENTER)
        _add_text(s, interval, x + 0.1, 2.05, 2.6, 0.3, size=16, color=BODY,
                  align=PP_ALIGN.CENTER)
        _add_text(s, label, x + 0.1, 2.45, 2.6, 0.3, size=14, color=MUTED,
                  align=PP_ALIGN.CENTER)

    _add_text(s, "Idea clave", 0.6, 3.3, 9.0, 0.3, size=18, bold=True, color=ACCENT)
    _add_multiline(s, [
        "El ancho del intervalo cambia segun la confianza del modelo en cada prediccion individual.",
        "Calibracion responde: \"cuanto riesgo hay?\". CP responde: \"con que confianza lo digo?\".",
        "En Mondrian, el cuantil se calcula por grado (A-G) para no ocultar subgrupos mal cubiertos.",
    ], 0.6, 3.7, 8.8, 1.5, size=16)

    # ════════════════════════════════════════════════════════
    # SLIDE 8: SECTION - METHODOLOGY
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _section_slide(s, "Seccion 2", "Metodologia")

    # ════════════════════════════════════════════════════════
    # SLIDE 9: METHODOLOGY 6 PHASES
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Seis fases secuenciales alineadas con CRISP-DM,\ndesde revision bibliografica hasta evaluacion regulatoria",
              0.5, 0.2, 9.0, 0.8, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    phases = [
        ("Fase 1", "Revision Bibliografica", "Scopus, arXiv, Google Scholar (2005-2026)"),
        ("Fase 2", "Construccion del Dataset", "1.86M prestamos, 110 cols, splits OOT"),
        ("Fase 3", "Modelado Base y Calibracion", "LR baseline + CatBoost + Venn-Abers"),
        ("Fase 4", "Conformal Prediction", "Mondrian PD, Adaptativo LGD, Split EAD"),
        ("Fase 5", "Evaluacion Comparativa", "Kupiec, Christoffersen, backtesting 35 meses"),
        ("Fase 6", "Impacto Regulatorio", "ECL escenarios, Stages IFRS 9, SICR conformal"),
    ]
    for i, (fase, titulo, detalle) in enumerate(phases):
        y_pos = 1.15 + i * 0.7
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos),
                                  Inches(1.2), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT
        box.line.fill.background()
        _add_text(s, fase, 0.55, y_pos + 0.1, 1.1, 0.35, size=13, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        _add_text(s, titulo, 1.85, y_pos + 0.02, 3.5, 0.3, size=16, bold=True, color=PRIMARY)
        _add_text(s, detalle, 1.85, y_pos + 0.3, 7.5, 0.25, size=13, color=MUTED)

    _add_text(s, "Cambio vs. anteproyecto: CatBoost como unico modelo (no XGBoost/LightGBM), seleccion automatica de calibrador, Mondrian como variante principal de CP.",
              0.5, 5.0, 9.0, 0.5, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 10: DATASET
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Lending Club: 1.86M prestamos con separacion temporal estricta\ny gradiente de default A(5.6%) a G(47.7%)",
              0.5, 0.2, 9.0, 0.8, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # Table: splits
    _add_table(s, ["Conjunto", "Registros", "Tasa Default", "Periodo"],
               [
                   ["Entrenamiento", "1,346,311", "18.52%", "2007-06 a 2017-03"],
                   ["Calibracion", "237,584", "22.20%", "2017-03 a 2017-12"],
                   ["Test (OOT)", "276,869", "21.98%", "2018-01 a 2020-09"],
               ], 0.5, 1.1, 5.5, font_size=13)

    # Image: default rate by grade
    img_path = NB_IMG / "01_eda_lending_club" / "default_rate_by_grade.png"
    _add_image_safe(s, img_path, 6.2, 1.0, 3.5, 2.3)

    _add_multiline(s, [
        "Separacion temporal OOT simula condiciones reales de produccion.",
        "El shift de default (18.5% -> 22%) incluye inicio de pandemia COVID-19.",
        "22 variables con data leakage eliminadas (total_pymnt, recoveries, etc.).",
        "42 features seleccionadas (9 categoricas, 27 numericas, 12 WOE).",
    ], 0.5, 3.3, 9.0, 1.8, size=16)

    # ════════════════════════════════════════════════════════
    # SLIDE 11: SECTION - RESULTS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _section_slide(s, "Seccion 3", "Resultados")

    # ════════════════════════════════════════════════════════
    # SLIDE 12: PD MODEL RESULTS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "CatBoost monotonico calibrado con Venn-Abers logra AUC 0.713\ny ECE 0.0067 — calibracion mejora calidad sin perder discriminacion",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # Table: model comparison
    _add_table(s, ["Modelo", "AUC-ROC", "Gini", "Brier", "D2-Brier", "ECE", "KS"],
               [
                   ["Reg. Logistica", "0.679", "0.358", "0.233", "-0.356", "—", "—"],
                   ["CatBoost (default)", "0.713", "0.426", "0.210", "-0.222", "—", "—"],
                   ["CatBoost (tuned)", "0.713", "0.426", "0.210", "-0.222", "—", "—"],
                   ["CatBoost (calibrado)", "0.713", "0.425", "0.155", "+0.099", "0.007", "0.313"],
               ], 0.5, 1.1, 5.5, font_size=11)

    # Image: ROC curves
    roc_path = NB_IMG / "03_pd_modeling" / "roc_curves.png"
    _add_image_safe(s, roc_path, 6.2, 1.0, 3.5, 2.2)

    _add_multiline(s, [
        "D2-Brier negativo = peor que predecir la prevalencia. Calibracion Venn-Abers lo corrige a +0.099.",
        "Restricciones monotonicas: annual_inc (neg.) y loan_to_income (pos.) — coherencia economica garantizada.",
        "Venn-Abers seleccionado automaticamente via politica multi-fold temporal (4 candidatos, 4 folds).",
        "ECE 0.0067 = buena calibracion, pero no responde: con que confianza lo digo? Eso lo resuelve CP.",
    ], 0.5, 3.3, 9.0, 1.8, size=15, bold_prefix=False)

    _add_text(s, "Evaluacion sobre 276,869 prestamos OOT (2018-01 a 2020-09).",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 13: CALIBRATION COMPARISON
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Venn-Abers gana la seleccion automatica multi-fold temporal\ncon el mejor ECE (0.00613) entre tres candidatos",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # Table: calibrator comparison
    _add_table(s, ["Calibrador", "AUC-ROC", "Brier Score", "ECE"],
               [
                   ["Platt Scaling", "0.7129", "0.1546", "0.00839"],
                   ["Isotonic Regression", "0.7128", "0.1545", "0.00625"],
                   ["Venn-Abers (ganador)", "0.7128", "0.1545", "0.00613"],
               ], 0.5, 1.1, 5.5, font_size=13)

    # Image: calibration curves
    cal_path = NB_IMG / "03_pd_modeling" / "calibration_curves.png"
    _add_image_safe(s, cal_path, 6.2, 1.0, 3.5, 2.2)

    _add_multiline(s, [
        "Seleccion automatica: 4 folds temporales, criterio Brier + AUC-guard (max -0.15pp).",
        "Venn-Abers produce dos probabilidades (p0, p1) que acotan la probabilidad real.",
        "Selector width promedio: 0.0005 — estimaciones muy estrechas (alta confianza).",
        "Cambio vs. anteproyecto: se paso de evaluacion manual a politica automatica objetiva.",
    ], 0.5, 3.3, 9.0, 1.8, size=16, bold_prefix=False)

    # ════════════════════════════════════════════════════════
    # SLIDE 14: CONFORMAL PD MONDRIAN
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Mondrian Conformal logra 92.42% de cobertura global\ncon 89.01% minima por grado al 90% de confianza",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # KPIs
    _kpi_box(s, 0.5, 1.15, "Coverage 90%", "92.42%", w=2.1)
    _kpi_box(s, 2.8, 1.15, "Coverage 95%", "95.93%", w=2.1)
    _kpi_box(s, 5.1, 1.15, "Min. por grado", "89.01%", w=2.1)
    _kpi_box(s, 7.4, 1.15, "Ancho medio", "0.7546", w=2.1)

    # Image: coverage by grade
    cov_path = NB_IMG / "04_conformal_prediction" / "coverage_by_grade.png"
    _add_image_safe(s, cov_path, 0.3, 2.3, 5.0, 2.5)

    _add_multiline(s, [
        "Todos los grados superan 87% de cobertura.",
        "Kupiec: p < 0.001 (esperado con N=276K).",
        "Christoffersen 90%: p_ind = 0.512 (violaciones independientes).",
        "Sobrecobertura deseable en contexto regulatorio: subestimar incertidumbre es peor.",
    ], 5.5, 2.4, 4.2, 2.5, size=15, bold_prefix=False)

    _add_text(s, "MAPIE 1.3.0 SplitConformalRegressor + ProbabilityRegressor wrapper por grado A-G.",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 15: MONDRIAN vs GLOBAL
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Mondrian supera a Split Global en +30pp de cobertura minima\npor grado con intervalos 21% mas eficientes",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # Comparison table
    _add_table(s, ["Metrica", "Global", "Mondrian", "Diferencia"],
               [
                   ["Cobertura 90%", "89.84%", "92.42%", "+2.58 pp"],
                   ["Min. cobertura por grado", "58.69%", "89.01%", "+30.32 pp"],
                   ["Ancho medio 90%", "0.955", "0.7546", "-21.3%"],
               ], 0.5, 1.1, 5.5, font_size=14)

    # Image: variant scatter or coverage-width tradeoff
    # Prefer publication heatmap; fall back to tradeoff plot
    heatmap_path = FIG_DIR / "p3_fig2_grade_coverage_heatmap.png"
    tradeoff_path = NB_IMG / "04_conformal_prediction" / "coverage_width_tradeoff.png"
    if not _add_image_safe(s, heatmap_path, 6.2, 1.0, 3.5, 2.2):
        _add_image_safe(s, tradeoff_path, 6.2, 1.0, 3.5, 2.2)

    _add_multiline(s, [
        "Global: 41.3% de grado G fuera del intervalo — invalida la garantia.",
        "Mondrian calibra quantiles por grado: precisos en A-B, amplios en F-G.",
        "Paradoja: mejor cobertura + intervalos mas estrechos (no hay trade-off).",
        "Implicacion operativa: decisiones por segmento necesitan garantia por segmento.",
    ], 0.5, 3.3, 9.0, 1.8, size=16, bold_prefix=False)

    # ════════════════════════════════════════════════════════
    # SLIDE 16: LGD / EAD RESULTS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "La triada PD-LGD-EAD tiene cobertura conformal verificada\n— extension a LGD y EAD cierra vacio en la literatura",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # LGD benchmark table
    _add_text(s, "Benchmark de 4 variantes LGD (60,850 defaults)", 0.5, 1.1, 5.5, 0.3,
              size=16, bold=True, color=ACCENT)
    _add_table(s, ["Variante LGD", "Cob. 90%", "Min Grado", "Ancho"],
               [
                   ["Two-stage split", "78.19%", "53.56%", "0.568"],
                   ["Direct split", "78.64%", "54.12%", "0.568"],
                   ["Direct CQR", "74.52%", "70.51%", "0.540"],
                   ["Adaptive grade-temporal", "90.50%", "90.47%", "0.496"],
               ], 0.5, 1.5, 5.5, font_size=12)

    # EAD box
    _kpi_box(s, 6.5, 1.2, "EAD Cob. 90%", "91.20%", w=2.5)
    _kpi_box(s, 6.5, 2.3, "EAD Cob. 95%", "95.28%", w=2.5)

    _add_multiline(s, [
        "Solo adaptive_grade_temporal paso todos los guardrails de cobertura, sesgo y eficiencia.",
        "LGD: variable continua en [0,1] con alta heteroscedasticidad por grado.",
        "EAD: alta precision porque la exposicion es contractual (ancho medio $132.58 USD).",
        "Contribucion: la mayoria de literatura (Bellotti, 2017) solo cubre PD, no LGD/EAD.",
    ], 0.5, 3.6, 9.0, 1.6, size=16, bold_prefix=False)

    # ════════════════════════════════════════════════════════
    # SLIDE 16b: THREE COMPLEMENTARY LAYERS (NEW)
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Tres capas complementarias de calidad del modelo\n— cada una resuelve un problema distinto",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    from pptx.enum.shapes import MSO_SHAPE as _MSO
    layers = [
        ("Capa 1: Calibracion", "Venn-Abers",
         ["Corrige nivel probabilistico", "ECE: 0.210 -> 0.0067",
          "D2-Brier: -0.222 -> +0.099"],
         "Cuanto riesgo hay?", ACCENT),
        ("Capa 2: Monotonia", "CatBoost constraints",
         ["Coherencia economica", "annual_inc: negativa",
          "loan_to_income: positiva", "0% violaciones"],
         "Tiene sentido economico?", GREEN),
        ("Capa 3: Conformal", "Mondrian CP",
         ["Cuantifica incertidumbre", "Cobertura: 92.42%",
          "Min grado: 89.01%", "Ancho: 0.7546"],
         "Con que confianza lo digo?", RGBColor(0xE6, 0x7E, 0x22)),
    ]
    for i, (title, method, bullets, question, col) in enumerate(layers):
        x = 0.3 + i * 3.2
        box = s.shapes.add_shape(_MSO.ROUNDED_RECTANGLE, Inches(x), Inches(1.15),
                                  Inches(3.0), Inches(3.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = col
        box.line.width = Pt(2)
        _add_text(s, title, x + 0.15, 1.2, 2.7, 0.3, size=15, bold=True, color=col,
                  align=PP_ALIGN.CENTER)
        _add_text(s, method, x + 0.15, 1.5, 2.7, 0.3, size=13, color=MUTED,
                  align=PP_ALIGN.CENTER)
        _add_multiline(s, bullets, x + 0.25, 1.85, 2.5, 1.6, size=14, color=BODY,
                       spacing=Pt(4))
        _add_text(s, f'"{question}"', x + 0.15, 3.65, 2.7, 0.5, size=13, bold=True,
                  color=col, align=PP_ALIGN.CENTER)

    _add_text(s, "Resultado: un modelo que no solo predice bien, sino que predice con coherencia y con incertidumbre cuantificada.",
              0.5, 4.6, 9.0, 0.5, size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    _add_text(s, "Ninguna capa por si sola es suficiente. Las tres juntas producen predicciones utiles para decision, no solo para ranking.",
              0.5, 5.15, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 17: IFRS9 ECL IMPACT
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Los intervalos conformales revelan $798M de incertidumbre\nen provisiones ECL — impacto material en lectura regulatoria",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # ECL scenarios table
    _add_table(s, ["Escenario", "ECL (USD)", "Uplift vs Base", "PD utilizada"],
               [
                   ["Baseline (PD puntual)", "$1,001M", "—", "PD calibrada"],
                   ["Mild stress", "$1,200M", "+22.8%", "PD + 0.5 x ancho CP"],
                   ["Adverse", "$1,463M", "+49.7%", "PD + 0.75 x ancho CP"],
                   ["Severe (PD alta CP)", "$1,799M", "+79.7%", "PD_high (95%)"],
               ], 0.5, 1.1, 6.0, font_size=13)

    # KPIs
    _kpi_box(s, 6.8, 1.1, "Stage 1", "34.51%", w=2.5)
    _kpi_box(s, 6.8, 2.15, "Stage 2 (SICR)", "43.01%", w=2.5)
    _kpi_box(s, 6.8, 3.2, "Stage 3 (impago)", "22.48%", w=2.5)

    # IFRS9 fan chart image
    fan_path = NB_IMG / "05_time_series_forecasting" / "ifrs9_scenario_fan_chart.png"
    _add_image_safe(s, fan_path, 0.3, 3.3, 5.0, 1.8)

    _add_multiline(s, [
        "Sin CP: el banco ve $1,001M y no sabe si es robusto o fragil.",
        "Con CP: rango $1,001M-$1,799M (+79.7%) — cuantificacion explicita.",
        "Implicacion: provisiones deben cubrir incertidumbre, no solo la media.",
        "Innovacion SICR: ancho conformal como senal de deterioro crediticio.",
    ], 5.5, 3.6, 4.2, 1.5, size=14, bold_prefix=True)

    _add_text(s, "De estimaciones puntuales a intervalos con garantia formal: la incertidumbre deja de ser invisible y se convierte en informacion accionable.",
              0.5, 5.15, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 18: SECTION - EVALUATOR & OBJECTIVES
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _section_slide(s, "Seccion 4", "Evaluacion y Cumplimiento")

    # ════════════════════════════════════════════════════════
    # SLIDE 19: EVALUATOR RESPONSE
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Evaluaciones aprobadas: anteproyecto 4.5/5.0 y documento 4.2/5.0\n— todas las correcciones solicitadas fueron implementadas",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    _add_text(s, "Evaluacion del anteproyecto (Feb 2026)", 0.5, 1.05, 5.5, 0.3,
              size=14, bold=True, color=ACCENT)
    _add_table(s, ["Criterio", "Nota", "Como se abordo"],
               [
                   ["Pertinencia", "4.8", "Coberturas >90%, $798M incertidumbre"],
                   ["Planteamiento", "4.3", "6 OE reformulados, medibles"],
                   ["Metodologia", "4.6", "CatBoost + Mondrian focalizado"],
                   ["Innovacion", "4.7", "CP triada PD-LGD-EAD + SICR"],
                   ["Viabilidad", "3.5", "Profundidad > amplitud"],
                   ["Comunicacion", "4.8", "Kupiec/Christoffersen + tablas"],
               ], 0.3, 1.35, 5.5, row_h=0.33, font_size=10)

    _add_text(s, "Evaluacion del documento (Mar 2026): 4.2/5.0", 6.0, 1.05, 3.7, 0.3,
              size=14, bold=True, color=ACCENT)
    _add_multiline(s, [
        "Correcciones solicitadas (todas implementadas):",
        "(a) Agregar resumen/abstract y palabras clave.",
        "(b) Reescribir conclusiones por objetivo (Tabla 12).",
        "(c) Estandarizar formato de citacion IEEE.",
    ], 6.2, 1.4, 3.5, 1.8, size=14, bold_prefix=False)

    _add_text(s, "Evaluador: Andres Felipe Garcia Ospina. Concepto documento: Aprobado con Ajustes (4.2/5.0, 19 de marzo de 2026).",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 20: METHODOLOGICAL CHANGES
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Los cambios metodologicos respecto al anteproyecto\nfueron informados por resultados experimentales",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    _add_table(s, ["Aspecto", "Anteproyecto", "Final", "Razon"],
               [
                   ["Modelo PD", "XGBoost", "CatBoost 1.2.8", "NaN y categoricas nativos"],
                   ["Modelos", "LR+XGB+LGB+CB", "LR + CatBoost", "Profundidad > amplitud"],
                   ["Calibrador", "Manual", "Auto multi-fold", "Venn-Abers gano"],
                   ["LGD", "Split unico", "4 variantes + guardrails", "Heteroscedasticidad"],
                   ["CP PD", "Split Global", "Mondrian A-G", "Min cov 58%->89%"],
                   ["Libreria", "MAPIE (s/v)", "MAPIE 1.3.0", "SplitConformalRegressor"],
               ], 0.3, 1.1, 9.4, row_h=0.38, font_size=11)

    _add_text(s, "Todos los cambios documentados con justificacion en el documento tecnico (Tabla 4).",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 21: OBJECTIVES FULFILLMENT
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Los seis objetivos especificos se cumplieron con evidencia\ncuantitativa y entregables verificables",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    _add_table(s, ["Objetivo", "Evidencia clave"],
               [
                   ["OE1: Marco teorico CP", "19 fuentes, 5 subsecciones, vacio PD-only identificado"],
                   ["OE2: Dataset OOT", "1.35M/237K/277K, 22 vars leakage, 42 features"],
                   ["OE3: PD calibrado", "AUC 0.713, ECE 0.007, Brier 0.155, KS 0.313"],
                   ["OE4: CP Mondrian PD/LGD/EAD", "PD 92.42%, LGD 90.50%, EAD 91.20% cob."],
                   ["OE5: Impacto IFRS 9", "ECL $1,001M-$1,799M (+79.7%), SICR conformal"],
                   ["OE6: Lineamientos adopcion", "Backtesting 35 meses, gate 7/13, dashboard"],
               ], 0.3, 1.1, 9.4, row_h=0.4, font_size=12)

    _add_multiline(s, [
        "Resultado en CRISP-DM: tesis trazable de punta a punta, no coleccion de experimentos aislados.",
        "Pipeline reproducible: DVC + MLflow + Streamlit (31 paginas de dashboard interactivo).",
        "690 tests automatizados pasando (0 fallos, 0 skips).",
    ], 0.5, 4.1, 9.0, 1.2, size=15, bold_prefix=False)

    # ════════════════════════════════════════════════════════
    # SLIDE 22: FEATURE IMPORTANCE (SHAP) + MONOTONICITY
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Las cinco variables mas importantes son consistentes con la\nliteratura de riesgo crediticio y refuerzan la interpretabilidad",
              0.5, 0.2, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    # SHAP image
    shap_path = NB_IMG / "03_pd_modeling" / "shap_summary.png"
    if not shap_path.exists():
        shap_path = NB_IMG / "13_model_explainability" / "shap_bar_importance.png"
    _add_image_safe(s, shap_path, 0.3, 1.1, 5.5, 3.5)

    _add_multiline(s, [
        "int_rate: mayor tasa = mayor riesgo.",
        "grade/sub_grade: gradiente confirmado.",
        "term: 60m = mayor riesgo que 36m.",
        "fico_score: factor protector.",
        "dti: indicador de fragilidad.",
    ], 6.0, 1.2, 3.7, 1.8, size=14, bold_prefix=True)

    # Monotonicity verification
    mono_path = NB_IMG / "13_model_explainability" / "monotonicity_verification.png"
    _add_image_safe(s, mono_path, 6.0, 3.2, 3.5, 1.5)

    _add_text(s, "SHAP + Auditoria de monotonia: 0% violaciones en annual_inc y loan_to_income.",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 24: CONCLUSIONS
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _add_text(s, "Conclusiones", 0.5, 0.2, 9.0, 0.4, size=20, color=RGBColor(0xA0, 0xBB, 0xDD))
    _add_rule(s, 0.5, 0.6, 9.0, 0.04, ACCENT)

    conclusions = [
        "1. Calibracion y CP son complementarios: Venn-Abers corrigio D2-Brier de -0.222 a +0.099. CP agrega bandas de incertidumbre con garantia formal de cobertura.",
        "",
        "2. Mondrian supera a Split Global: +30pp en cobertura minima por grado (89% vs 59%) con intervalos 21% mas eficientes. La garantia solo promedio oculta fallos graves.",
        "",
        "3. Triada PD-LGD-EAD con CP: PD 92.42%, LGD 90.50%, EAD 91.20%. Cierra vacio en literatura (Bellotti, 2017) que solo cubria PD.",
        "",
        "4. Impacto IFRS 9: ECL de $1,001M a $1,799M (+79.7%) — $798M de incertidumbre cuantificada. Ancho conformal como senal SICR.",
        "",
        "5. Restricciones monotonicas garantizan coherencia economica sin sacrificio discriminativo (AUC invariante: 0.7127 vs 0.7128).",
    ]
    _add_multiline(s, conclusions, 0.5, 0.8, 9.0, 3.8, size=18, color=WHITE, spacing=Pt(4))

    _add_text(s, "carlos.vergara@utp.edu.co  |  Repositorio: github.com/cavr/lending-club-risk-project",
              0.5, 4.8, 7.0, 0.4, size=14, color=RGBColor(0xA0, 0xBB, 0xDD))

    # ════════════════════════════════════════════════════════
    # SLIDE 25: FUTURE WORK
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Lineas futuras: de la especializacion a la maestria\ny la adopcion en entornos bancarios reales",
              0.5, 0.2, 9.0, 0.8, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.0, 9.0)

    futures = [
        "CP Adaptativa (ACI): Ajustar quantiles dinamicamente ante cambios en la distribucion (no estacionariedad financiera).",
        "Optimizacion robusta de portafolio: Usar intervalos conformales como conjuntos de incertidumbre en formulaciones de Pyomo/HiGHS.",
        "CP online / streaming: Pipeline de decisiones de credito en tiempo real con actualizaciones incrementales.",
        "Extension IFRS 9: Validar ancho conformal como senal primaria de SICR en portafolios bancarios reales.",
        "Validacion multi-mercado: Aplicar framework a hipotecas, tarjetas y prestamos corporativos.",
    ]
    _add_multiline(s, futures, 0.6, 1.1, 8.8, 3.5, size=17, bold_prefix=True, spacing=Pt(12))

    _add_text(s, "Las tres primeras lineas corresponden al plan de maestria (ya en desarrollo).",
              0.5, 5.1, 9.0, 0.35, size=12, color=MUTED)

    # ════════════════════════════════════════════════════════
    # SLIDE 26: REFERENCES
    # ════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_text(s, "Referencias", 0.5, 0.2, 9.0, 0.5, size=24, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 0.7, 9.0)

    refs = [
        "[1] Angelopoulos & Bates (2023). A Gentle Introduction to Conformal Prediction. FnTML, 16(4), 494-591.",
        "[2] Baesens, Roesch & Scheule (2016). Credit Risk Analytics. Wiley.",
        "[3] Basel Committee (2006). International Convergence of Capital Measurement.",
        "[4] Bellotti (2017). Reliable region predictions for automated credit scoring. COPA, PMLR.",
        "[5] Deloitte (2020). IFRS 9 and Expected Credit Loss: Modelling Challenges.",
        "[6] Fontana, Zeni & Vantini (2023). Conformal Prediction: Unified Review. Bernoulli, 29(1).",
        "[7] IASB (2014). IFRS 9 Financial Instruments.",
        "[8] Javanmardi & Vovk (2023). Multi probability predictions with Venn-Abers. COPA, PMLR.",
        "[9] Lessmann et al. (2015). Benchmarking classification for credit scoring. EJOR, 247(1).",
        "[10] Niculescu-Mizil & Caruana (2005). Predicting good probabilities. ICML.",
        "[11] Platt (1999). Probabilistic Outputs for SVMs. MIT Press.",
        "[12] Romano, Patterson & Candes (2019). Conformalized Quantile Regression. NeurIPS.",
        "[13] Vovk, Gammerman & Shafer (2022). Algorithmic Learning in a Random World, 2nd ed. Springer.",
        "[14] Vovk & Petej (2014). Venn-Abers Predictors. UAI.",
        "[15] MAPIE Contributors (2023). MAPIE: Model Agnostic Prediction Interval Estimator.",
    ]
    _add_multiline(s, refs, 0.5, 0.85, 9.0, 4.5, size=12, color=BODY, spacing=Pt(3))

    # ════════════════════════════════════════════════════════
    # APPENDIX SLIDES
    # ════════════════════════════════════════════════════════

    # APPENDIX A: Statistical calibration tests
    s = prs.slides.add_slide(blank)
    _add_text(s, "Apendice A — Tests estadisticos de calibracion conformal", 0.5, 0.15, 9.0, 0.35,
              size=14, color=MUTED)
    _add_text(s, "Tests formales confirman cobertura valida con independencia\nde violaciones en el tiempo",
              0.5, 0.5, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.25, 9.0)

    _add_table(s, ["Test", "Nivel 90%", "Nivel 95%", "Interpretacion"],
               [
                   ["Kupiec p-value", "< 0.001", "< 0.001", "Rechaza cobertura exacta (N=276K)"],
                   ["Christoffersen p_ind", "0.512", "0.034", "Violaciones independientes (90%)"],
                   ["KS p-value", "0.000427", "—", "Sensibilidad a desviaciones finas"],
                   ["Kuiper p-value", "0.000153", "—", "Confirma calibracion precisa"],
                   ["Spiegelhalter p-value", "1.28e-06", "—", "Alta sensibilidad estadistica"],
               ], 0.3, 1.4, 9.4, font_size=12)

    _add_multiline(s, [
        "Con N=276K, incluso desviaciones de 2.42pp son significativas. Lo relevante: cobertura supera el objetivo.",
        "Christoffersen confirma que las violaciones no se agrupan temporalmente (p=0.512 > 0.05).",
    ], 0.5, 4.0, 9.0, 1.0, size=15)

    # APPENDIX B: Tech Stack
    s = prs.slides.add_slide(blank)
    _add_text(s, "Apendice B — Stack tecnologico y reproducibilidad", 0.5, 0.15, 9.0, 0.35,
              size=14, color=MUTED)
    _add_text(s, "Pipeline completamente reproducible con 690 tests,\nDVC, MLflow y dashboard Streamlit de 31 paginas",
              0.5, 0.5, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.25, 9.0)

    _add_table(s, ["Categoria", "Herramientas"],
               [
                   ["ML / Modelado", "CatBoost 1.2.8, scikit-learn 1.6.1, Optuna 4.7"],
                   ["Conformal", "MAPIE 1.3.0 (SplitConformalRegressor, Mondrian)"],
                   ["Evaluacion", "Kupiec, Christoffersen, Winkler Score"],
                   ["Datos", "pandas, numpy, pandera (validacion de esquemas)"],
                   ["Visualizacion", "matplotlib, seaborn, SHAP 0.48"],
                   ["Desarrollo", "Python 3.12, uv, ruff, pytest 690 tests, loguru"],
                   ["MLOps", "DVC, MLflow, Streamlit (31 paginas)"],
               ], 0.3, 1.4, 9.4, font_size=12)

    _add_multiline(s, [
        "Todo el codigo fuente, notebooks y artefactos estan documentados para reproduccion completa.",
        "Pre-commit: ruff lint+format, nbstripout, trailing-whitespace, check-yaml/toml.",
    ], 0.5, 4.3, 9.0, 0.8, size=15)

    # APPENDIX C: Monte Carlo ECL
    s = prs.slides.add_slide(blank)
    _add_text(s, "Apendice C — Simulacion Monte Carlo de ECL con 16,384 escenarios", 0.5, 0.15, 9.0, 0.35,
              size=14, color=MUTED)
    _add_text(s, "Distribucion de cola del ECL confirma que el escenario severo\nes consistente con la simulacion estocastica",
              0.5, 0.5, 9.0, 0.8, size=22, bold=True, color=PRIMARY)
    _add_rule(s, 0.5, 1.25, 9.0)

    _add_table(s, ["Estadistico", "Valor"],
               [
                   ["Media ECL", "$1,059.8M"],
                   ["P50 (Mediana)", "$997.3M"],
                   ["P90", "$1,571.5M"],
                   ["P95", "$1,748.6M"],
                   ["P99", "$2,111.9M"],
                   ["Expected Shortfall (ES@95%)", "$1,975.3M"],
               ], 0.5, 1.4, 5.0, font_size=13)

    _add_multiline(s, [
        "16,384 escenarios simulados con seed=42.",
        "ES@95%: $1,975M — perdida esperada dado que se supera P95.",
        "Rango intercuartilico: $762M de incertidumbre intrinseca.",
        "GPU speedup vs CPU: 75.58x.",
    ], 6.0, 1.5, 3.7, 2.5, size=15, bold_prefix=False)

    # ═══════════ SAVE ═══════════
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
